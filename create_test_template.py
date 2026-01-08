#!/usr/bin/env python3
"""Create a test PPT template with :TAG: format placeholders."""

from pptx import Presentation
from pptx.util import Inches, Pt

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Slide 1: Title with tags
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

# Title box
title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(1))
title_frame = title_box.text_frame
title_frame.text = ":TITLE:"
title_frame.paragraphs[0].font.size = Pt(44)
title_frame.paragraphs[0].font.bold = True

# Company info
company_box = slide1.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
company_frame = company_box.text_frame
company_frame.text = "Company: :COMPANY_NAME:\nSponsor: :SPONSOR_NAME:"
company_frame.paragraphs[0].font.size = Pt(24)

# Location
location_box = slide1.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
location_frame = location_box.text_frame
location_frame.text = "Location: :LOCATION_COUNTRY:, :LOCATION_STATE:"
location_frame.paragraphs[0].font.size = Pt(20)

# Slide 2: Project details with mixed formats
slide2 = prs.slides.add_slide(prs.slide_layouts[6])

# Header
header_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
header_frame = header_box.text_frame
header_frame.text = "Project Overview"
header_frame.paragraphs[0].font.size = Pt(36)
header_frame.paragraphs[0].font.bold = True

# Details using tags
details_box = slide2.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(3))
details_frame = details_box.text_frame
details_text = """Asset Type: :ASSET_TYPE:
Technology: :TECHNOLOGY:
Capacity: :UNIT:
Status: :PROJECT_STATUS:
COD: :COD:

Investment: :INITIAL_INVESTMENT: (Initial) + :FUTURE_INVESTMENT: (Future)
Contractor: :CONTRACTOR:
Offtaker: :OFFTAKER:"""
details_frame.text = details_text
details_frame.paragraphs[0].font.size = Pt(18)

# Slide 3: AI-generated content with tags
slide3 = prs.slides.add_slide(prs.slide_layouts[6])

# Title
ai_title_box = slide3.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
ai_title_frame = ai_title_box.text_frame
ai_title_frame.text = "Project Description"
ai_title_frame.paragraphs[0].font.size = Pt(36)
ai_title_frame.paragraphs[0].font.bold = True

# AI-generated teaser using tags
ai_box = slide3.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4.5))
ai_frame = ai_box.text_frame
ai_frame.text = "[AI: Write a professional 3-paragraph description about {COMPANY_NAME}'s {TECHNOLOGY} project in {LOCATION_COUNTRY}. The project has a capacity of {UNIT} and requires an initial investment of {INITIAL_INVESTMENT}. Follow this structure: Sentence 1: introduce the project. Sentence 2: describe the technology. Sentence 3: explain the investment and timeline with COD of {COD}.]"
ai_frame.paragraphs[0].font.size = Pt(16)
ai_frame.word_wrap = True

# Slide 4: Simple description field
slide4 = prs.slides.add_slide(prs.slide_layouts[6])

desc_title_box = slide4.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
desc_title_frame = desc_title_box.text_frame
desc_title_frame.text = "Technical Overview"
desc_title_frame.paragraphs[0].font.size = Pt(36)
desc_title_frame.paragraphs[0].font.bold = True

desc_box = slide4.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4))
desc_frame = desc_box.text_frame
desc_frame.text = ":DESCRIPTION:"
desc_frame.paragraphs[0].font.size = Pt(20)
desc_frame.word_wrap = True

# Save
prs.save('test_template_tags.pptx')
print("✅ Created test_template_tags.pptx with :TAG: placeholders")
print("\n📋 Template includes:")
print("  - Slide 1: Title and basic info with :COMPANY_NAME:, :SPONSOR_NAME:, etc.")
print("  - Slide 2: Project details with multiple tags")
print("  - Slide 3: AI-generated content using {TAG} in prompts")
print("  - Slide 4: Description field")
