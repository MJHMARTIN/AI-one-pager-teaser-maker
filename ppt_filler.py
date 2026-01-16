"""
PowerPoint Filler Module
Handles placeholder detection and replacement in PPTX files.
"""

import re
import logging
import pandas as pd
from pptx import Presentation
from pptx.presentation import Presentation as PresentationType
from ai_generators import generate_beta_text
from field_mapping import create_placeholder_mapping
from excel_loader import normalize_label

# Setup logging
logger = logging.getLogger(__name__)


def resolve_tag_to_value(
    tag: str,
    data_dict: dict,
    placeholder_mapping: dict,
    module_type: str = None
) -> str:
    """
    Resolve a :TAG: placeholder to its value using the mapping and data dictionary.
    
    Args:
        tag: The tag name (e.g., 'COMPANY_NAME')
        data_dict: Dictionary of normalized_label -> value
        placeholder_mapping: Dictionary of TAG -> [possible label variations]
        module_type: Detected module type ("Module1", "Module2", "Module3")
    
    Returns:
        The resolved value or empty string if not found
    """
    from field_mapping import get_module_field_mappings, fuzzy_match_label
    
    logger.info(f"🔍 resolve_tag_to_value called for: {tag}")
    logger.info(f"   Module type: {module_type}")
    logger.info(f"   data_dict size: {len(data_dict)}")
    logger.info(f"   data_dict keys sample: {list(data_dict.keys())[:5]}")
    
    # If we have a detected module, prioritize module-specific labels
    if module_type:
        module_mappings = get_module_field_mappings()
        if module_type in module_mappings and tag in module_mappings[module_type]:
            module_label = module_mappings[module_type][tag]
            normalized = normalize_label(module_label)
            logger.info(f"   Module mapping: {tag} → '{module_label}' (normalized: '{normalized}')")
            if normalized in data_dict:
                value = data_dict[normalized]
                logger.info(f"   ✅ FOUND in data_dict: '{value[:50]}'")
                return value
            else:
                logger.warning(f"   ❌ NOT FOUND: '{normalized}' not in data_dict keys")
                logger.warning(f"   Available keys: {list(data_dict.keys())[:10]}")
    
    # Check if this tag is in our mapping
    if tag in placeholder_mapping:
        possible_labels = placeholder_mapping[tag]
        
        # Try each possible label variation
        for label_variation in possible_labels:
            normalized = normalize_label(label_variation)
            if normalized in data_dict:
                logger.debug(f"Resolved :TAG: {tag} using label: {label_variation}")
                return data_dict[normalized]
        
        # If no exact match, try fuzzy matching
        fuzzy_match = fuzzy_match_label(possible_labels[0], list(data_dict.keys()))
        if fuzzy_match and fuzzy_match in data_dict:
            logger.debug(f"Resolved :TAG: {tag} using fuzzy match: {fuzzy_match}")
            return data_dict[fuzzy_match]
    
    # Try using the tag itself as a label (for backwards compatibility)
    normalized_tag = normalize_label(tag)
    if normalized_tag in data_dict:
        logger.debug(f"Resolved :TAG: {tag} using normalized tag")
        return data_dict[normalized_tag]
    
    # Try fuzzy match on the tag
    from field_mapping import fuzzy_match_label
    fuzzy_match = fuzzy_match_label(tag, list(data_dict.keys()))
    if fuzzy_match and fuzzy_match in data_dict:
        logger.debug(f"Resolved :TAG: {tag} using fuzzy match on tag")
        return data_dict[fuzzy_match]
    
    logger.warning(f"Could not resolve :TAG: {tag}")
    return ""


def process_placeholder(
    placeholder: str,
    row_data: pd.Series,
    excel_columns: list,
    beta_tone: str,
    missing_to_blank: bool,
    data_dict: dict = None,
    placeholder_mapping: dict = None,
    excel_format: str = 'column-based',
    module_type: str = None,
    log_warnings: bool = True
) -> str:
    """
    Resolve a placeholder to final text using smart local generation.

    Supports multiple formats:
    - Direct: "Title" → row_data["Title"] (column-based)
    - Tag: ":COMPANY_NAME:" → resolved via mapping (row-based)
    - AI: "AI: Write teaser about {Company}" → smart text generation
    
    Args:
        placeholder: The placeholder text to process
        row_data: DataFrame row (for column-based mode)
        excel_columns: List of available Excel columns
        beta_tone: "short", "medium", or "long"
        missing_to_blank: If True, return empty string for missing fields
        data_dict: Dictionary of normalized labels to values (for hybrid mode)
        placeholder_mapping: Mapping of tags to field names
        excel_format: 'column-based', 'hybrid', or 'row-based'
        module_type: Detected module type for specialized mappings
        log_warnings: Whether to log warnings for missing fields
    
    Returns:
        The resolved text
    """
    logger.debug(f"Processing placeholder: {placeholder[:80]}...")
    
    # Handle :TAG: format (new row-based system)
    tag_match = re.match(r'^:([A-Z_]+):$', placeholder.strip())
    if tag_match and data_dict is not None and placeholder_mapping is not None:
        tag = tag_match.group(1)
        logger.info(f"🔍 PROCESSING TAG: {tag}")
        logger.info(f"   data_dict has {len(data_dict)} keys")
        logger.info(f"   placeholder_mapping has entry for {tag}: {tag in placeholder_mapping}")
        result = resolve_tag_to_value(tag, data_dict, placeholder_mapping, module_type)
        if result:
            logger.info(f"✅ Resolved :TAG: {tag} → {result[:50]}...")
        else:
            logger.warning(f"❌ Failed to resolve :TAG: {tag} - returned empty string")
        return result
    
    # AI placeholder - always use local generation
    if placeholder.startswith("AI:"):
        logger.info(f"Processing AI placeholder: {placeholder[:80]}...")
        prompt_text = placeholder[3:].strip()
        
        # Handle double "AI:" at the start (common error)
        if prompt_text.startswith("AI:"):
            prompt_text = prompt_text[3:].strip()

        # Replace {{ColumnName}} or {ColumnName} with actual values from Excel row
        for col in excel_columns:
            if col in row_data.index:
                value = "" if pd.isna(row_data[col]) else str(row_data[col])
                # Support both {{}} and {} formats
                prompt_text = prompt_text.replace(f"{{{{{col}}}}}", value)  # {{ColumnName}}
                prompt_text = prompt_text.replace(f"{{{col}}}", value)       # {ColumnName}
        
        # Also support {TAG} format for row-based data
        if data_dict is not None and placeholder_mapping is not None:
            # Find all {TAG} or {{TAG}} patterns
            tokens = re.findall(r'\{\{?([A-Z_]+)\}?\}?', prompt_text)
            for token in tokens:
                value = resolve_tag_to_value(token, data_dict, placeholder_mapping, module_type)
                prompt_text = prompt_text.replace(f"{{{{{token}}}}}", value)
                prompt_text = prompt_text.replace(f"{{{token}}}", value)
        
        # Fix common typos: {{TOKEN} with only one closing brace
        prompt_text = re.sub(r'\{\{([^}]+)\}(?!\})', r'{{\1}}', prompt_text)
        
        # Find remaining unreplaced tokens
        remaining_tokens = re.findall(r'\{\{([^}]+)\}\}', prompt_text)
        remaining_tokens.extend(re.findall(r'\{([^}]+)\}', prompt_text))
        
        # CRITICAL: If any tokens couldn't be replaced, DO NOT call AI!
        if remaining_tokens:
            missing_fields = ', '.join(set(remaining_tokens))
            error_msg = f"[CANNOT GENERATE: Missing Excel data for {missing_fields}]"
            if log_warnings:
                logger.error(f"AI generation blocked - Missing Excel columns in prompt: {missing_fields}")
            return error_msg if not missing_to_blank else ""
        
        # All tokens successfully replaced - safe to call AI
        logger.debug(f"All tokens replaced successfully. Calling AI with: {prompt_text[:200]}...")
        
        generated = generate_beta_text(prompt_text, row_data, beta_tone)
        
        logger.info(f"✓ AI generated text: {generated[:100]}...")
        return generated

    # Direct placeholder - check if we have row_data (column-based) or data_dict (hybrid/row-based)
    if data_dict is not None and placeholder_mapping is not None:
        # Using hybrid/row-based mode - try to resolve from data_dict
        col_name = placeholder.strip()
        normalized = normalize_label(col_name)
        
        # First try: direct lookup in data_dict
        if normalized in data_dict:
            logger.debug(f"Resolved direct placeholder '{col_name}' from data_dict")
            return data_dict[normalized]
        
        # Second try: Check if this is a TAG name (e.g., "INDUSTRY") and use module mapping
        if col_name.isupper() and col_name in placeholder_mapping:
            # Try to resolve using TAG mapping system
            tag_value = resolve_tag_to_value(col_name, data_dict, placeholder_mapping, module_type)
            if tag_value:
                logger.info(f"Resolved direct placeholder '{col_name}' via TAG mapping: {tag_value[:50]}...")
                return tag_value
        
        # Third try: fuzzy match
        from field_mapping import fuzzy_match_label
        fuzzy_match = fuzzy_match_label(col_name, list(data_dict.keys()))
        if fuzzy_match and fuzzy_match in data_dict:
            logger.info(f"Resolved '{col_name}' via fuzzy match: {fuzzy_match}")
            return data_dict[fuzzy_match]
        
        if log_warnings:
            logger.warning(f"Missing field: {col_name}")
        return "" if missing_to_blank else f"[MISSING FIELD: {col_name}]"
    
    # Legacy column-based mode with row_data
    col_name = placeholder.strip()
    if row_data is not None and col_name in row_data.index:
        val = row_data[col_name]
        result = "" if pd.isna(val) else str(val)
        logger.debug(f"Resolved column '{col_name}' from row_data")
        return result
    else:
        if log_warnings:
            logger.warning(f"Missing column: {col_name}")
        return "" if missing_to_blank else f"[MISSING COLUMN: {col_name}]"


def extract_placeholders_from_text(text: str) -> dict:
    """
    Extract all placeholders from text.
    
    Returns:
        dict with keys 'ai' (AI placeholders), 'direct' (direct placeholders), 'tags' (TAG placeholders)
    """
    # AI placeholders: [AI: ...] (handle multi-line)
    ai_matches = re.findall(r"\[AI:\s*(.+?)\]", text, re.DOTALL)
    
    # If no matches, try without closing bracket (for malformed placeholders)
    if not ai_matches:
        ai_no_close = re.findall(r"\[AI:\s*(.+)$", text, re.DOTALL)
        if ai_no_close:
            logger.warning("Found AI placeholder without closing bracket ]")
            ai_matches = ai_no_close
    
    # Direct placeholders: [...] (but not [AI: ...])
    direct_matches = re.findall(r"\[(?!AI:)(.+?)\]", text)
    
    # TAG placeholders: :TAG:
    tag_matches = re.findall(r":([A-Z_]+):", text)
    tags = [f":{tag}:" for tag in tag_matches]
    
    logger.debug(f"Extracted {len(ai_matches)} AI, {len(direct_matches)} direct, {len(tags)} TAG placeholders")
    
    return {
        'ai': ai_matches,
        'direct': direct_matches,
        'tags': tags
    }


def replace_text_preserving_format(text_frame, old_text: str, new_text: str):
    """
    Replace text in a text frame while preserving formatting.
    
    Args:
        text_frame: The pptx text_frame object
        old_text: Original text
        new_text: New text to insert
    """
    if not text_frame.paragraphs:
        # No paragraphs, just set text directly
        logger.debug("No paragraphs in text_frame, setting text directly")
        text_frame.text = new_text
        return
    
    para = text_frame.paragraphs[0]
    if not para.runs:
        # No runs, just set text
        logger.debug("No runs in paragraph, setting text directly")
        para.text = new_text
        return
    
    # Get formatting from first run
    run0 = para.runs[0]
    font = run0.font
    
    # Clear and recreate
    text_frame.clear()
    new_para = text_frame.paragraphs[0]
    new_run = new_para.add_run()
    new_run.text = new_text
    
    # Copy formatting
    if font.size:
        new_run.font.size = font.size
    if font.name:
        new_run.font.name = font.name
    if font.bold is not None:
        new_run.font.bold = font.bold
    if font.italic is not None:
        new_run.font.italic = font.italic
    if font.color and hasattr(font.color, 'rgb'):
        new_run.font.color.rgb = font.color.rgb
    
    # Copy paragraph formatting
    new_para.alignment = para.alignment
    new_para.level = para.level
    
    logger.debug("Replaced text with format preservation")


def iter_shapes(shapes):
    """
    Recursively iterate through all shapes, including those in groups.
    """
    for shape in shapes:
        yield shape
        # If shape is a group, recursively process shapes inside it
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            try:
                for subshape in iter_shapes(shape.shapes):
                    yield subshape
            except AttributeError:
                pass


def process_shape_text(shape, row_data, excel_columns, beta_tone, missing_to_blank, 
                       data_dict, placeholder_mapping, excel_format, module_type, stats):
    """
    Process text in a single shape and replace placeholders.
    Returns True if text was modified.
    """
    if not hasattr(shape, "text_frame"):
        return False
    
    original_text = shape.text
    if not original_text:
        return False
    
    logger.debug(f"Processing shape: {shape.name} (text: {original_text[:50]}...)")
    
    placeholders = extract_placeholders_from_text(original_text)
    
    if not placeholders['ai'] and not placeholders['direct'] and not placeholders['tags']:
        return False
    
    new_text = original_text
    
    # Process direct placeholders
    for placeholder in placeholders['direct']:
        replacement = process_placeholder(
            placeholder, row_data, excel_columns, beta_tone,
            missing_to_blank, data_dict, placeholder_mapping,
            excel_format, module_type, log_warnings=True
        )
        new_text = new_text.replace(f"[{placeholder}]", replacement)
        stats['direct_replaced'] += 1
        logger.debug(f"Replaced [{placeholder}] with: {replacement[:50]}...")
    
    # Process TAG placeholders
    for tag in placeholders['tags']:
        replacement = process_placeholder(
            tag, row_data, excel_columns, beta_tone,
            missing_to_blank, data_dict, placeholder_mapping,
            excel_format, module_type, log_warnings=True
        )
        new_text = new_text.replace(tag, replacement)
        stats['tags_replaced'] += 1
        logger.debug(f"Replaced {tag} with: {replacement[:50]}...")
    
    # Process AI placeholders
    for ai_prompt in placeholders['ai']:
        logger.info(f"Generating AI text for: {ai_prompt[:60]}...")
        replacement = process_placeholder(
            f"AI: {ai_prompt}", row_data, excel_columns, beta_tone,
            missing_to_blank, data_dict, placeholder_mapping,
            excel_format, module_type, log_warnings=False
        )
        new_text = new_text.replace(f"[AI: {ai_prompt}]", replacement)
        stats['ai_generated'] += 1
        logger.info(f"AI generated: {replacement[:80]}...")
    
    # Apply changes if text was modified
    if new_text != original_text:
        replace_text_preserving_format(shape.text_frame, original_text, new_text)
        stats['processed'] += 1
        logger.debug(f"Updated shape: {shape.name}")
        return True
    
    return False


def fill_pptx(
    prs: PresentationType,
    row_data: pd.Series,
    excel_columns: list,
    beta_tone: str,
    missing_to_blank: bool,
    data_dict: dict = None,
    placeholder_mapping: dict = None,
    excel_format: str = 'column-based',
    module_type: str = None,
    progress_callback=None
) -> dict:
    """
    Fill all placeholders in a PowerPoint presentation.
    
    Args:
        prs: The python-pptx Presentation object
        row_data: DataFrame row (for column-based mode)
        excel_columns: List of available Excel columns
        beta_tone: "short", "medium", or "long"
        missing_to_blank: If True, return empty string for missing fields
        data_dict: Dictionary of normalized labels to values (for hybrid mode)
        placeholder_mapping: Mapping of tags to field names
        excel_format: 'column-based', 'hybrid', or 'row-based'
        module_type: Detected module type for specialized mappings
        progress_callback: Optional function(current, total) to report progress
    
    Returns:
        dict with statistics: {'total_shapes': int, 'processed': int, 'ai_generated': int, 'direct_replaced': int}
    """
    logger.info("=== Starting PPT fill ===")
    logger.info(f"Mode: {excel_format}, Module: {module_type or 'auto'}, Tone: {beta_tone}")
    
    # CRITICAL DEBUG: Log what we received
    logger.info("="*80)
    logger.info("RECEIVED IN fill_pptx:")
    logger.info(f"  data_dict: {len(data_dict) if data_dict else 0} keys")
    if data_dict:
        logger.info(f"  First 10 keys: {list(data_dict.keys())[:10]}")
        for key in list(data_dict.keys())[:3]:
            logger.info(f"    '{key}' = '{data_dict[key][:50] if data_dict[key] else '(empty)'}'")
    logger.info(f"  placeholder_mapping: {len(placeholder_mapping) if placeholder_mapping else 0} tags")
    if placeholder_mapping:
        logger.info(f"  First 5 tags: {list(placeholder_mapping.keys())[:5]}")
    logger.info("="*80)
    
    stats = {
        'total_shapes': 0,
        'processed': 0,
        'ai_generated': 0,
        'direct_replaced': 0,
        'tags_replaced': 0
    }
    
    # Count all text targets (including grouped shapes and table cells)
    total_targets = 0
    for s in prs.slides:
        for shp in iter_shapes(s.shapes):
            if hasattr(shp, "text_frame"):
                total_targets += 1
            if getattr(shp, "has_table", False):
                tbl = shp.table
                for r in tbl.rows:
                    total_targets += len(r.cells)
    
    stats['total_shapes'] = total_targets
    logger.info(f"Found {total_targets} text targets across {len(prs.slides)} slides")
    
    processed = 0
    
    for slide_idx, slide in enumerate(prs.slides):
        logger.debug(f"Processing slide {slide_idx + 1}/{len(prs.slides)}")
        
        # Use recursive iterator to get ALL shapes including grouped ones
        for shape in iter_shapes(slide.shapes):
            
            # Handle tables
            if getattr(shape, "has_table", False):
                logger.debug(f"Processing table: {shape.name}")
                table = shape.table
                
                for row_idx, r in enumerate(table.rows):
                    for cell_idx, cell in enumerate(r.cells):
                        original_text = cell.text
                        if not original_text:
                            processed += 1
                            if progress_callback:
                                progress_callback(processed, total_targets)
                            continue
                        
                        placeholders = extract_placeholders_from_text(original_text)
                        
                        if not placeholders['ai'] and not placeholders['direct'] and not placeholders['tags']:
                            processed += 1
                            if progress_callback:
                                progress_callback(processed, total_targets)
                            continue
                        
                        new_text = original_text
                        
                        # Process direct placeholders
                        for placeholder in placeholders['direct']:
                            replacement = process_placeholder(
                                placeholder, row_data, excel_columns, beta_tone,
                                missing_to_blank, data_dict, placeholder_mapping,
                                excel_format, module_type, log_warnings=True
                            )
                            new_text = new_text.replace(f"[{placeholder}]", replacement)
                            stats['direct_replaced'] += 1
                            logger.debug(f"Replaced [{placeholder}] with: {replacement[:50]}...")
                        
                        # Process TAG placeholders
                        for tag in placeholders['tags']:
                            logger.info(f"📍 TABLE CELL [{row_idx},{cell_idx}]: Processing TAG {tag}")
                            logger.info(f"   Original text: {original_text[:100]}")
                            replacement = process_placeholder(
                                tag, row_data, excel_columns, beta_tone,
                                missing_to_blank, data_dict, placeholder_mapping,
                                excel_format, module_type, log_warnings=True
                            )
                            logger.info(f"   Replacement value: '{replacement[:100] if replacement else '(EMPTY)'}'")
                            new_text = new_text.replace(tag, replacement)
                            stats['tags_replaced'] += 1
                            logger.info(f"✓ Replaced {tag} with: {replacement[:50]}...")
                        
                        # Process AI placeholders
                        for ai_prompt in placeholders['ai']:
                            logger.info(f"Generating AI text for: {ai_prompt[:60]}...")
                            replacement = process_placeholder(
                                f"AI: {ai_prompt}", row_data, excel_columns, beta_tone,
                                missing_to_blank, data_dict, placeholder_mapping,
                                excel_format, module_type, log_warnings=False
                            )
                            new_text = new_text.replace(f"[AI: {ai_prompt}]", replacement)
                            stats['ai_generated'] += 1
                        
                        # Apply changes if text was modified
                        if new_text != original_text:
                            replace_text_preserving_format(cell.text_frame, original_text, new_text)
                            stats['processed'] += 1
                            logger.debug(f"Updated table cell [{row_idx},{cell_idx}]")
                        
                        processed += 1
                        if progress_callback:
                            progress_callback(processed, total_targets)
                
                continue  # Move to next shape
            
            # Handle regular text shapes (including grouped shapes)
            if process_shape_text(shape, row_data, excel_columns, beta_tone, missing_to_blank,
                                 data_dict, placeholder_mapping, excel_format, module_type, stats):
                logger.debug(f"Processed shape: {shape.name}")
            
            processed += 1
            if progress_callback:
                progress_callback(processed, total_targets)
    
    logger.info("=== PPT fill complete ===")
    logger.info(f"Stats: {stats}")
    
    return stats
