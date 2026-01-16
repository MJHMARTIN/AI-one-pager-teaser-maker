"""
PPTX Teaser Generator - Streamlit UI
Refactored to use modular architecture with separate concerns.
"""

import streamlit as st
import pandas as pd
from pptx import Presentation
from io import BytesIO
import logging

# Import modules
from excel_loader import detect_excel_format, parse_hybrid_excel, parse_multi_sheet_excel
from module_detector import detect_module_type
from field_mapping import create_placeholder_mapping
from ppt_filler import fill_pptx

# Setup logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# ---------------- Streamlit page config ----------------
st.set_page_config(page_title="PPTX Teaser Generator with AI", layout="wide")
st.title("📊 PPTX Teaser Generator with AI")

logger.info("Application started")

st.write(
    "Upload your template PPTX and Excel data. "
    "The template can contain direct placeholders like [Title] and AI prompts like "
    "[AI: Write a short teaser about {Company} in {Industry}]. "
    "Powered by Perplexity API for intelligent content generation. "
    "Formatting (font, size, color, bold, italic) of each text box is preserved."
)

# ---------------- Sidebar: Config ----------------
st.sidebar.header("⚙️ Configuration")

st.sidebar.success("🚀 **Powered by Perplexity API** — Advanced AI text generation with intelligent fallbacks!")

beta_tone = st.sidebar.radio(
    "AI Text Tone",
    options=["short", "medium", "long"],
    index=1,
    help="short: 1 sentence. medium: 1-2 sentences. long: 2-3 detailed sentences.",
)

missing_to_blank = st.sidebar.checkbox(
    "Treat missing columns as blank",
    value=True,
    help="If a placeholder column is missing, insert blank instead of an error tag.",
)

st.sidebar.divider()

st.sidebar.info(
    """
**Template syntax**

**1. Direct placeholders (no AI)**  
- **Column-based Excel:**  
  In PowerPoint: `Title: [Title]` or `Company: [Company Name]`  
  In Excel: columns `Title`, `Company Name`

- **Row-based Excel:**  
  In PowerPoint: `:COMPANY_NAME:`, `:TITLE:`, `:SPONSOR_NAME:`  
  In Excel: rows with labels like "Company name", "Title", "Sponsor name"

**2. AI-generated prompts (Perplexity API)**  
- In PowerPoint:  
  `[AI: Write a teaser about {Company} in {Industry}]`  
  or  
  `[AI: Write about {COMPANY_NAME} in {LOCATION_COUNTRY}]`  
- Column-based: tokens match Excel columns  
- Row-based: tokens match canonical tags (COMPANY_NAME, etc.)  
- Uses Perplexity's LLM for intelligent, context-aware generation

**Row-based format** automatically handles label variations:  
- "Company name", "Company Name", "Name of company" → all work!  
- Uses fuzzy matching for flexibility

Text box formatting is preserved when content is replaced.
"""
)

# ---------------- File uploads ----------------
col1, col2 = st.columns(2)

with col1:
    st.subheader("Step 1: Template PPTX")
    template_file = st.file_uploader(
        "Upload your PPTX template (with [Placeholders] and [AI: ...] prompts)",
        type="pptx",
        key="template",
    )

with col2:
    st.subheader("Step 2: Excel Data")
    excel_file = st.file_uploader(
        "Upload your Excel file (columns must match placeholders)",
        type=["xlsx", "xlsm"],
        key="excel",
        help="Supports .xlsx and .xlsm files. Macros in .xlsm files are ignored for security."
    )

# ---------------- Template inspector (optional helper) ----------------
if template_file:
    with st.expander("🔍 Inspect template text boxes (optional)"):
        try:
            prs_preview = Presentation(template_file)
            st.write(f"Slides: **{len(prs_preview.slides)}**")
            slide_index = st.selectbox(
                "Slide to inspect",
                options=range(len(prs_preview.slides)),
                format_func=lambda i: f"Slide {i+1}",
            )
            slide = prs_preview.slides[slide_index]
            st.write(f"**Shapes on Slide {slide_index + 1}:**")

            for idx, shape in enumerate(slide.shapes):
                # Show text from text shapes and table cells
                if hasattr(shape, "text_frame"):
                    st.write(f"- Shape {idx} name: `{shape.name}`")
                    st.write(f"  Text: `{shape.text[:80]}`")
                elif getattr(shape, "has_table", False):
                    st.write(f"- Table {idx} with {len(shape.table.rows)} rows")
        except Exception as e:
            st.error(f"Error inspecting template: {e}")

# ---------------- Main generation flow ----------------
if template_file and excel_file:
    st.divider()

    # Read Excel - try multi-sheet first
    try:
        # Detect file type and add info for .xlsm
        file_name = excel_file.name if hasattr(excel_file, 'name') else 'unknown'
        if file_name.endswith('.xlsm'):
            st.info("ℹ️ **XLSM file detected**: Reading data only, macros are disabled for security")
        
        logger.info(f"Reading Excel file: {file_name}")
        
        # Check if multi-sheet (pandas reads .xlsm same as .xlsx, ignores macros)
        all_sheets = pd.read_excel(excel_file, sheet_name=None, engine='openpyxl')
        num_sheets = len(all_sheets)
        
        # For single sheet, use traditional flow
        if num_sheets == 1:
            df = list(all_sheets.values())[0]
            is_multi_sheet = False
        else:
            # Multi-sheet Excel
            is_multi_sheet = True
            st.info(f"📚 **Multi-sheet Excel detected:** {num_sheets} sheets found")
            logger.info(f"Multi-sheet mode: {num_sheets} sheets")
            
            # Show sheet names
            with st.expander("📄 Sheet names"):
                for sheet_name in all_sheets.keys():
                    st.write(f"- {sheet_name}")
    except Exception as e:
        st.error(f"Error reading Excel: {e}")
        logger.error(f"Excel read error: {e}")
        st.stop()

    # Handle single vs multi-sheet
    if not is_multi_sheet:
        if df.empty:
            st.error("Excel file has no rows.")
            st.stop()

        st.subheader("📋 Excel preview")
        st.dataframe(df, use_container_width=True)

        # Detect Excel format
        excel_format = detect_excel_format(df)
        st.info(f"📋 **Detected Excel format:** {excel_format}")
        logger.info(f"Detected format: {excel_format}")
    else:
        # Multi-sheet: show first sheet as preview
        first_sheet_name = list(all_sheets.keys())[0]
        first_sheet_df = all_sheets[first_sheet_name]
        
        st.subheader(f"📋 Excel preview (showing first sheet: '{first_sheet_name}')")
        st.dataframe(first_sheet_df, use_container_width=True)
        
        st.info(f"📋 **Hybrid mode**: Reading both column headers and row-based Label/Value pairs from all sheets")
    
    # Parse data based on sheet count
    if is_multi_sheet:
        # Use multi-sheet parser
        use_namespacing = st.checkbox(
            "Use sheet namespacing",
            value=False,
            help="Prefix field names with sheet name (e.g., 'financials.total_cost')"
        )
        
        data_dict, sheet_info = parse_multi_sheet_excel(excel_file, use_namespacing)
        placeholder_mapping = create_placeholder_mapping()
        
        st.success(f"✅ Parsed {len(data_dict)} total fields from {num_sheets} sheets")
        logger.info(f"Parsed {len(data_dict)} fields from multi-sheet Excel")
        
        with st.expander("🔍 View parsed data by sheet"):
            for info in sheet_info:
                st.write(f"**{info['name']}**: {info['format']} - {info['fields']} fields")
        
        with st.expander("� View ALL parsed Excel fields (comprehensive)"):
            st.write(f"**Total fields extracted: {len(data_dict)}**")
            st.write("Below are ALL labels and values found in your Excel file:")
            st.write("---")
            
            # Show all fields in a nice table
            if data_dict:
                import pandas as pd
                display_data = []
                for key, value in sorted(data_dict.items()):
                    value_preview = str(value)[:80] + "..." if len(str(value)) > 80 else str(value)
                    display_data.append({
                        "Excel Label (normalized)": key,
                        "Value Preview": value_preview
                    })
                
                df_display = pd.DataFrame(display_data)
                st.dataframe(df_display, use_container_width=True)
                
                st.write("---")
                st.info("💡 **Tip:** You can use any of these labels in your PPT template as `:LABEL:` format. The system uses fuzzy matching, so slight variations work!")
            else:
                st.warning("No fields were extracted. Check if your Excel has a Label/Value format.")
        
        # For multi-sheet, we don't need row selection
        row_data = None
        excel_columns = []
        excel_format = 'hybrid'  # Track that we're using hybrid mode
        
    else:
        # Single sheet: use hybrid parser (reads both column and row formats)
        data_dict = parse_hybrid_excel(df)
        placeholder_mapping = create_placeholder_mapping()
        excel_format = 'hybrid'  # Track that we're using hybrid mode
        
        st.info(f"📋 **Hybrid mode**: Reading both column headers and row-based Label/Value pairs")
        st.success(f"✅ Parsed {len(data_dict)} fields from Excel")
        logger.info(f"Parsed {len(data_dict)} fields in hybrid mode")
        
        with st.expander("� View ALL parsed Excel fields (comprehensive)"):
            st.write(f"**Total fields extracted: {len(data_dict)}**")
            st.write("Below are ALL labels and values found in your Excel file:")
            st.write("---")
            
            # Show all fields in a nice table
            if data_dict:
                display_data = []
                for key, value in sorted(data_dict.items()):
                    value_preview = str(value)[:80] + "..." if len(str(value)) > 80 else str(value)
                    display_data.append({
                        "Excel Label (normalized)": key,
                        "Value Preview": value_preview
                    })
                
                df_display = pd.DataFrame(display_data)
                st.dataframe(df_display, use_container_width=True)
                
                st.write("---")
                st.info("💡 **Tip:** You can use any of these labels in your PPT template as `:LABEL:` format. The system uses fuzzy matching, so slight variations work!")
            else:
                st.warning("No fields were extracted. Check if your Excel has a Label/Value format.")
            for key, value in list(data_dict.items())[:20]:
                st.write(f"**{key}**: {value[:100] if value else '(empty)'}...")
            if len(data_dict) > 20:
                st.write(f"... and {len(data_dict) - 20} more fields")
        
        # For hybrid, we don't need to select a row
        row_data = pd.Series(dtype='object')  # Empty series for compatibility
        excel_columns = []

    # Load PPTX template
    try:
        prs = Presentation(template_file)
        logger.info(f"Loaded PPTX with {len(prs.slides)} slides")
    except Exception as e:
        st.error(f"Error loading PPTX template: {e}")
        logger.error(f"PPTX load error: {e}")
        st.stop()

    st.divider()
    st.subheader("⚙️ Generation settings")

    # Info about hybrid mode
    if is_multi_sheet:
        st.info(f"ℹ️ Hybrid mode: Using unified data from all {num_sheets} sheets")
    else:
        st.info("ℹ️ Hybrid mode: Using both column headers and row-based data")

    # Button to generate
    if st.button("🚀 Generate PPTX", type="primary"):
        try:
            logger.info("=== Starting PPTX generation ===")
            
            # Detect module type from data
            module_type = detect_module_type(data_dict)
            if module_type:
                st.success(f"✅ Using {module_type} field mappings")
                logger.info(f"Detected module type: {module_type}")
            
            # CRITICAL DEBUG: Log the actual data_dict being passed to PPT filler
            logger.info("="*80)
            logger.info("CRITICAL DEBUG: data_dict being passed to fill_pptx:")
            logger.info(f"  Keys count: {len(data_dict)}")
            logger.info(f"  First 10 keys: {list(data_dict.keys())[:10]}")
            for key in list(data_dict.keys())[:5]:
                value_preview = str(data_dict[key])[:50]
                logger.info(f"  '{key}' = '{value_preview}'")
            logger.info("="*80)
            
            # Use the new fill_pptx function from ppt_filler module
            progress = st.progress(0.0)
            status = st.empty()
            
            def progress_callback(current, total):
                progress.progress(current / max(total, 1))
            
            stats = fill_pptx(
                prs=prs,
                row_data=row_data,
                excel_columns=excel_columns,
                beta_tone=beta_tone,
                missing_to_blank=missing_to_blank,
                data_dict=data_dict,
                placeholder_mapping=placeholder_mapping,
                excel_format=excel_format,
                module_type=module_type,
                progress_callback=progress_callback
            )
            
            progress.progress(1.0)
            status.text("✅ Done")
            
            logger.info(f"Generation complete: {stats}")
            st.info(f"📊 **Statistics**: "
                   f"{stats['ai_generated']} AI texts generated, "
                   f"{stats['direct_replaced']} direct replacements, "
                   f"{stats['tags_replaced']} tags replaced")

            # Save and offer download
            output = BytesIO()
            prs.save(output)
            output.seek(0)

            st.success("✅ PPTX generated successfully!")
            st.download_button(
                label="📥 Download filled PPTX",
                data=output.getvalue(),
                file_name="teaser_filled.pptx",
                mime=(
                    "application/"
                    "vnd.openxmlformats-officedocument."
                    "presentationml.presentation"
                ),
            )
            
            logger.info("PPTX download ready")

        except Exception as e:
            st.error(f"Error generating PPTX: {e}")
            logger.error(f"Generation error: {e}", exc_info=True)
            st.exception(e)

logger.info("UI render complete")
