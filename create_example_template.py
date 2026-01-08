#!/usr/bin/env python3
"""Create PPT template with user's example prompts."""

from pptx import Presentation
from pptx.util import Inches, Pt

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Slide 1: Example 1 - Battery Storage / Japan
slide1 = prs.slides.add_slide(prs.slide_layouts[6])

title1_box = slide1.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
title1_frame = title1_box.text_frame
title1_frame.text = "Example 1: Battery Storage / Japan"
title1_frame.paragraphs[0].font.size = Pt(32)
title1_frame.paragraphs[0].font.bold = True

ai1_box = slide1.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
ai1_frame = ai1_box.text_frame
ai1_frame.text = '''[AI: Write a 3-paragraph Project Highlight section for a {Sector} sector project in {Country} called "{Project_Focus}". 

Paragraph 1: What does the company do? Use: "{Company_Operations}"

Paragraph 2: What is the project scope and partnerships? Use: "{Service_Scope}". Also mention they "{Partnership}" and "{Partnership_Benefit}".

Paragraph 3: What is the investment plan? Use: "The project will commence with an initial investment of {Initial_Investment}, with a projected expansion to {Expansion_Investment} {Expansion_Path}.".]'''
ai1_frame.paragraphs[0].font.size = Pt(14)
ai1_frame.word_wrap = True

# Slide 2: Example 2 - Solar / Vietnam
slide2 = prs.slides.add_slide(prs.slide_layouts[6])

title2_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
title2_frame = title2_box.text_frame
title2_frame.text = "Example 2: Solar / Vietnam"
title2_frame.paragraphs[0].font.size = Pt(32)
title2_frame.paragraphs[0].font.bold = True

ai2_box = slide2.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
ai2_frame = ai2_box.text_frame
ai2_frame.text = '''[AI: Write a 3-paragraph Project Highlight for a {Sector} project in {Country} called "{Project_Focus}".

Paragraph 1: What does the company do? Use: "{Company_Operations}"

Paragraph 2: Describe the project scope and partnerships. Use: "{Service_Scope}". Also state they "{Partnership}" and "{Partnership_Benefit}".

Paragraph 3: What is the investment roadmap? Use: "The project will commence with an initial investment of {Initial_Investment}, with a projected expansion to {Expansion_Investment} {Expansion_Path}".]'''
ai2_frame.paragraphs[0].font.size = Pt(14)
ai2_frame.word_wrap = True

# Slide 3: Example 3 - Wind / India
slide3 = prs.slides.add_slide(prs.slide_layouts[6])

title3_box = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
title3_frame = title3_box.text_frame
title3_frame.text = "Example 3: Wind / India"
title3_frame.paragraphs[0].font.size = Pt(32)
title3_frame.paragraphs[0].font.bold = True

ai3_box = slide3.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
ai3_frame = ai3_box.text_frame
ai3_frame.text = '''[AI: Write a 3-paragraph Project Highlight for a {Sector} initiative in {Country} called "{Project_Focus}".

Paragraph 1: What is the company's business? Use: "{Company_Operations}"

Paragraph 2: Explain the project scope and key partnerships. Use: "{Service_Scope}". Also include that they "{Partnership}" and "{Partnership_Benefit}".

Paragraph 3: Outline the capital structure. Use: "The project will commence with an initial investment of {Initial_Investment}, with a projected expansion to {Expansion_Investment} {Expansion_Path}".]'''
ai3_frame.paragraphs[0].font.size = Pt(14)
ai3_frame.word_wrap = True

# Save
prs.save('test_template_examples.pptx')
print("✅ Created test_template_examples.pptx with user's example prompts")
