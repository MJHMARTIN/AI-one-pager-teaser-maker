"""
PPTX Teaser Generator - Streamlit UI
Refactored to use modular architecture with separate concerns.
"""

import streamlit as st
import pandas as pd
from pptx import Presentation
from io import BytesIO
import logging

# Import modules
from excel_loader import detect_excel_format, parse_hybrid_excel, parse_multi_sheet_excel
from module_detector import detect_module_type
from field_mapping import create_placeholder_mapping
from ppt_filler import fill_pptx

# Setup logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# ---------------- Streamlit page config ----------------
st.set_page_config(page_title="PPTX Teaser Generator with AI", layout="wide")
st.title("📊 PPTX Teaser Generator with AI")

logger.info("Application started")

st.write(
    "Upload your template PPTX and Excel data. "
    "The template can contain direct placeholders like [Title] and AI prompts like "
    "[AI: Write a short teaser about {Company} in {Industry}]. "
    "Formatting (font, size, color, bold, italic) of each text box is preserved."
)

# ---------------- Sidebar: Config ----------------
st.sidebar.header("⚙️ Configuration")

st.sidebar.info("💡 **Using smart local AI generator** — No API needed, completely free!")

beta_tone = st.sidebar.radio(
    "AI Text Tone",
    options=["short", "medium", "long"],
    index=1,
    help="short: 1 sentence. medium: 1-2 sentences. long: 2-3 detailed sentences.",
)

missing_to_blank = st.sidebar.checkbox(
    "Treat missing columns as blank",
    value=True,
    help="If a placeholder column is missing, insert blank instead of an error tag.",
)

st.sidebar.divider()

st.sidebar.info(
    """
**Template syntax**

**1. Direct placeholders (no AI)**  
- **Column-based Excel:**  
  In PowerPoint: `Title: [Title]` or `Company: [Company Name]`  
  In Excel: columns `Title`, `Company Name`

- **Row-based Excel:**  
  In PowerPoint: `:COMPANY_NAME:`, `:TITLE:`, `:SPONSOR_NAME:`  
  In Excel: rows with labels like "Company name", "Title", "Sponsor name"

**2. AI-generated prompts (local)**  
- In PowerPoint:  
  `[AI: Write a teaser about {Company} in {Industry}]`  
  or  
  `[AI: Write about {COMPANY_NAME} in {LOCATION_COUNTRY}]`  
- Column-based: tokens match Excel columns  
- Row-based: tokens match canonical tags (COMPANY_NAME, etc.)

**Row-based format** automatically handles label variations:  
- "Company name", "Company Name", "Name of company" → all work!  
- Uses fuzzy matching for flexibility

Text box formatting is preserved when content is replaced.
"""
)

# ---------------- File uploads ----------------
col1, col2 = st.columns(2)

with col1:
    st.subheader("Step 1: Template PPTX")
    template_file = st.file_uploader(
        "Upload your PPTX template (with [Placeholders] and [AI: ...] prompts)",
        type="pptx",
        key="template",
    )

with col2:
    st.subheader("Step 2: Excel Data")
    excel_file = st.file_uploader(
        "Upload your Excel file (columns must match placeholders)",
        type=["xlsx", "xlsm"],
        key="excel",
        help="Supports .xlsx and .xlsm files. Macros in .xlsm files are ignored for security."
    )
# ---------------- File uploads ----------------
    """Normalize Excel label for consistent matching."""
    if not isinstance(label, str):
        return ""
    # Lowercase, strip spaces, remove punctuation
    normalized = label.lower().strip()
    # Remove common punctuation but keep underscores
    normalized = re.sub(r'[^\w\s]', '', normalized)
    # Replace multiple spaces with single space
    normalized = re.sub(r'\s+', ' ', normalized)
    return normalized


def detect_excel_format(df: pd.DataFrame) -> str:
    """
    Detect if Excel is column-based or row-based format.
    
    Returns:
        'column-based': Traditional format where each column is a field
        'row-based': Label-Value pair format where rows contain field definitions
    """
    # Check if first two columns look like Label-Value pairs
    if len(df.columns) >= 2:
        first_col = df.columns[0]
        second_col = df.columns[1]
        
        # Common indicators of row-based format
        label_indicators = ['label', 'field', 'name', 'key', 'question', 'item']
        value_indicators = ['value', 'answer', 'data', 'content', 'response']
        
        first_lower = str(first_col).lower()
        second_lower = str(second_col).lower()
        
        # Check column names
        is_label_value = (
            any(ind in first_lower for ind in label_indicators) and
            any(ind in second_lower for ind in value_indicators)
        )
        
        if is_label_value:
            return 'row-based'
        
        # Check content: if first column has question-like text
        if len(df) > 0:
            first_values = df.iloc[:, 0].dropna().astype(str).head(5)
            # Check if values look like labels (contain spaces, question marks, colons)
            label_like_count = sum(
                1 for v in first_values 
                if len(v) > 10 and (' ' in v or ':' in v or '?' in v)
            )
            if label_like_count >= 3:  # If 3+ rows look like labels
                return 'row-based'
    
    # Default to column-based
    return 'column-based'


def parse_hybrid_excel(df: pd.DataFrame, sheet_name: str = None, use_namespacing: bool = False) -> dict:
    """
    Parse Excel in BOTH column-based and row-based formats simultaneously.
    
    - Column-based: First row (headers) become field names, second row has values
    - Row-based: First column = labels, second column = values
    
    This hybrid approach handles any Excel structure.
    
    Args:
        df: DataFrame to parse
        sheet_name: Name of the sheet (for namespacing)
        use_namespacing: If True, prefix keys with sheet name
    
    Returns:
        dict: {normalized_label: value}
    """
    data_dict = {}
    
    # PART 1: Read as COLUMN-based (headers = field names)
    if len(df) > 0 and len(df.columns) > 0:
        # Use first data row as values
        first_row = df.iloc[0]
        for col_name in df.columns:
            if pd.isna(col_name) or str(col_name).strip() == '':
                continue
            
            normalized = normalize_label(str(col_name))
            if use_namespacing and sheet_name:
                sheet_prefix = normalize_label(sheet_name)
                normalized = f"{sheet_prefix}.{normalized}"
            
            value = first_row[col_name]
            data_dict[normalized] = "" if pd.isna(value) else str(value)
    
    # PART 2: Read as ROW-based (column 0 = labels, column 1 = values)
    if len(df.columns) >= 2:
        label_col = df.columns[0]
        value_col = df.columns[1]
        
        for idx, row in df.iterrows():
            label = row[label_col]
            value = row[value_col]
            
            # Skip empty labels
            if pd.isna(label) or str(label).strip() == '':
                continue
            
            normalized = normalize_label(str(label))
            if use_namespacing and sheet_name:
                sheet_prefix = normalize_label(sheet_name)
                normalized = f"{sheet_prefix}.{normalized}"
            
            # Only add if not already present from column-based reading
            # (row-based takes precedence if there's overlap)
            data_dict[normalized] = "" if pd.isna(value) else str(value)
    
    return data_dict


def parse_row_based_excel(df: pd.DataFrame, sheet_name: str = None, use_namespacing: bool = False) -> dict:
    """
    Parse row-based Excel format into a normalized key-value dictionary.
    
    Expects Excel with structure:
    - Column 0: Labels/Questions
    - Column 1: Values/Answers
    - Optional: Column 2+ for categories/sections
    
    Args:
        df: DataFrame to parse
        sheet_name: Name of the sheet (for namespacing)
        use_namespacing: If True, prefix keys with sheet name (e.g., "financials.total_cost")
    
    Returns:
        dict: {normalized_label: value}
    """
    data_dict = {}
    
    # Use first two columns as label-value
    label_col = df.columns[0]
    value_col = df.columns[1] if len(df.columns) > 1 else df.columns[0]
    
    for idx, row in df.iterrows():
        label = row[label_col]
        value = row[value_col]
        
        # Skip empty labels
        if pd.isna(label) or str(label).strip() == '':
            continue
        
        # Normalize the label
        normalized = normalize_label(str(label))
        
        # Add sheet namespacing if requested
        if use_namespacing and sheet_name:
            sheet_prefix = normalize_label(sheet_name)
            normalized = f"{sheet_prefix}.{normalized}"
        
        # Store the value (convert NaN to empty string)
        data_dict[normalized] = "" if pd.isna(value) else str(value)
    
    return data_dict


def parse_multi_sheet_excel(excel_file, use_namespacing: bool = False) -> tuple:
    """
    Read all sheets from an Excel file and merge into unified knowledge base.
    
    Args:
        excel_file: File-like object or path to Excel (.xlsx or .xlsm)
        use_namespacing: If True, prefix keys with sheet names
    
    Returns:
        tuple: (unified_data_dict, sheet_info_list)
    """
    # Read all sheets (supports both .xlsx and .xlsm, macros ignored)
    sheets_dict = pd.read_excel(excel_file, sheet_name=None, engine='openpyxl')
    
    unified_data = {}
    sheet_info = []
    
    for sheet_name, df in sheets_dict.items():
        if df.empty:
            continue
        
        # In multi-sheet mode, try to parse everything as HYBRID (both row and column)
        # This reads both column headers as fields AND label/value pairs
        try:
            sheet_data = parse_hybrid_excel(df, sheet_name, use_namespacing)
            
            if sheet_data:  # Only count if we got data
                # Merge into unified dict (later sheets override earlier ones if no namespacing)
                unified_data.update(sheet_data)
                
                sheet_info.append({
                    'name': sheet_name,
                    'format': 'hybrid (row+column)',
                    'fields': len(sheet_data)
                })
            else:
                sheet_info.append({
                    'name': sheet_name,
                    'format': 'empty or invalid',
                    'fields': 0
                })
        except Exception as e:
            sheet_info.append({
                'name': sheet_name,
                'format': f'error: {str(e)[:30]}',
                'fields': 0
            })
    
    return unified_data, sheet_info


def fuzzy_match_label(target: str, available_labels: list, threshold: float = 0.6) -> str:
    """
    Find the best matching label using fuzzy string matching.
    
    Args:
        target: The label to match
        available_labels: List of available labels
        threshold: Minimum similarity score (0-1)
    
    Returns:
        Best matching label or None
    """
    normalized_target = normalize_label(target)
    
    # Try exact match first
    if normalized_target in available_labels:
        return normalized_target
    
    # Try fuzzy matching
    matches = get_close_matches(normalized_target, available_labels, n=1, cutoff=threshold)
    
    if matches:
        return matches[0]
    
    return None


def get_module_field_mappings() -> dict:
    """
    Define module-specific field mappings for different Excel formats.
    
    Each module (1, 2, 3) uses different labels for the same logical field.
    This maps each module's specific labels to canonical PPT placeholder tags.
    
    Returns:
        dict: {module_name: {PPT_TAG: excel_label}}
    """
    return {
        "Module1": {
            "ISSUER": "sponsor name",
            "INDUSTRY": "primary focus area",
            "JURISDICTION": "country of incorporation",
            "ISSUANCE_TYPE": "financing type",
            "INITIAL_NOTIONAL": "total financing amount",
            "COUPON_RATE": "coupon rate",
            "COUPON_FREQUENCY": "coupon frequency",
            "TENOR": "requested tenor",
            "CLIENT_SUMMARY": "sponsor summary",
            "PROJECT_HIGHLIGHT": "sponsor background investment strategy"
        },
        "Module2": {
            "ISSUER": "company legal name",
            "INDUSTRY": "primary industry",
            "JURISDICTION": "country of incorporation",
            "ISSUANCE_TYPE": "financing type",
            "INITIAL_NOTIONAL": "total financing amount",
            "COUPON_RATE": "coupon rate",
            "COUPON_FREQUENCY": "coupon frequency",
            "TENOR": "requested tenor",
            "CLIENT_SUMMARY": "company overview business model",
            "PROJECT_HIGHLIGHT": "company growth strategy financial projections"
        },
        "Module3": {
            "ISSUER": "project name",
            "INDUSTRY": "project type",
            "JURISDICTION": "project location country",
            "ISSUANCE_TYPE": "financing type",
            "INITIAL_NOTIONAL": "total project cost",
            "COUPON_RATE": "coupon rate",
            "COUPON_FREQUENCY": "coupon frequency",
            "TENOR": "project tenor",
            "CLIENT_SUMMARY": "project description",
            "PROJECT_HIGHLIGHT": "project overview technical specs impact"
        }
    }


def detect_module_type(data_dict: dict) -> str:
    """
    Detect which module type based on Excel field names.
    
    Args:
        data_dict: Dictionary of normalized_label -> value
    
    Returns:
        Module type: "Module1", "Module2", "Module3", or None
    """
    module_mappings = get_module_field_mappings()
    
    # Score each module based on how many of its specific fields are present
    scores = {}
    
    for module_name, field_map in module_mappings.items():
        score = 0
        # Check unique fields for each module
        for tag, excel_label in field_map.items():
            normalized = normalize_label(excel_label)
            if normalized in data_dict:
                score += 1
        scores[module_name] = score
    
    # Return module with highest score (if any fields matched)
    if scores:
        best_module = max(scores, key=scores.get)
        if scores[best_module] > 0:
            st.info(f"🔍 **Detected format:** {best_module}")
            return best_module
    
    return None


def create_placeholder_mapping() -> dict:
    """
    Define the mapping between PPT placeholder tags and canonical field names.
    
    This allows flexible wording in Excel while maintaining consistent PPT templates.
    
    Returns:
        dict: {PPT_TAG: [list of possible Excel label variations]}
    """
    base_mapping = {
        'COMPANY_NAME': ['company name', 'company', 'name of company', 'issuer', 'issuer name'],
        'SPONSOR_NAME': ['sponsor name', 'sponsor', 'lead sponsor', 'project sponsor'],
        'LOCATION_COUNTRY': ['location country', 'country', 'jurisdiction', 'location'],
        'LOCATION_STATE': ['location state', 'state', 'province', 'region'],
        'ASSET_TYPE': ['asset type', 'type of asset', 'asset class', 'sector'],
        'PROJECT_TYPE': ['project type', 'type of project', 'development type'],
        'UNIT': ['unit', 'capacity', 'size', 'scale'],
        'TECHNOLOGY': ['technology', 'tech', 'technology type', 'technical approach'],
        'INITIAL_INVESTMENT': ['initial investment', 'initial capex', 'phase 1 investment', 'startup cost'],
        'FUTURE_INVESTMENT': ['future investment', 'expansion investment', 'phase 2 investment', 'growth capex'],
        'CONTRACTOR': ['contractor', 'epc contractor', 'construction partner', 'builder'],
        'OFFTAKER': ['offtaker', 'offtake partner', 'purchaser', 'buyer'],
        'PROJECT_STATUS': ['project status', 'status', 'development stage', 'phase'],
        'COD': ['cod', 'commercial operation date', 'operational date', 'completion date'],
        'DESCRIPTION': ['description', 'project description', 'overview', 'summary'],
        'TITLE': ['title', 'project title', 'project name', 'deal name'],
        'INDUSTRY': ['industry', 'sector', 'vertical', 'market'],
    }
    
    # Add module-specific mappings
    module_mappings = get_module_field_mappings()
    for module_name, field_map in module_mappings.items():
        for tag, excel_label in field_map.items():
            normalized = normalize_label(excel_label)
            if tag in base_mapping:
                # Add to existing list if not already there
                if normalized not in base_mapping[tag]:
                    base_mapping[tag].append(normalized)
            else:
                # Create new entry
                base_mapping[tag] = [normalized]
    
    return base_mapping


# ---------------- Smart local text generation ----------------
# ---------------- FIXED TEMPLATE FUNCTIONS ----------------

def generate_sector_label(prompt_text: str, max_words: int = 3) -> str:
    """
    Generate a 2-3 word sector label.
    
    Args:
        prompt_text: The prompt with substituted values
        max_words: Maximum number of words (default 3)
    
    Returns:
        A concise sector label (e.g., "Solar Energy", "Wind Energy")
    """
    prompt_lower = prompt_text.lower()
    
    # Keyword-based sector detection
    sector_patterns = {
        'solar': 'Solar Energy',
        'wind': 'Wind Energy',
        'battery': 'Energy Storage',
        'storage': 'Energy Storage',
        'hydro': 'Hydro Energy',
        'nuclear': 'Nuclear Energy',
        'geothermal': 'Geothermal Energy',
        'biomass': 'Biomass Energy',
        'hydrogen': 'Hydrogen Energy',
        'data center': 'Data Centers',
        'infrastructure': 'Infrastructure',
        'real estate': 'Real Estate',
        'technology': 'Technology',
        'healthcare': 'Healthcare',
        'manufacturing': 'Manufacturing',
    }
    
    # Check for keywords in prompt
    for keyword, label in sector_patterns.items():
        if keyword in prompt_lower:
            return label
    
    # Check for "new energy" or "renewable"
    if 'new energy' in prompt_lower or 'renewable' in prompt_lower:
        return 'New Energy'
    
    # Extract quoted values as fallback
    quoted_values = re.findall(r'["\']([^"\']+)["\']', prompt_text)
    if quoted_values:
        # Take first quoted value and limit to max_words
        words = quoted_values[0].split()[:max_words]
        return ' '.join(words).title()
    
    return 'New Energy'


def generate_client_oneliner(prompt_text: str) -> str:
    """
    Generate a one-line anonymous client description.
    Pattern: "The client is [developing/operating] {asset_description} in {country}."
    
    Args:
        prompt_text: The prompt with substituted values
    
    Returns:
        A single sentence describing the client anonymously
    """
    prompt_lower = prompt_text.lower()
    
    # Extract quoted values (asset description and country)
    # First try longer strings (likely asset descriptions)
    quoted_values = re.findall(r'["\']([^"\']{20,})["\']', prompt_text)
    
    # If we didn't find long strings, try shorter ones
    if not quoted_values:
        quoted_values = re.findall(r'["\']([^"\']{5,})["\']', prompt_text)
    
    if len(quoted_values) >= 2:
        asset_description = quoted_values[0]
        country = quoted_values[1]
        
        # Determine action verb based on prompt content
        if 'operat' in prompt_lower and 'develop' in prompt_lower:
            verb_phrase = 'developing and operating'
        elif 'operat' in prompt_lower:
            verb_phrase = 'operating'
        else:
            verb_phrase = 'developing'
        
        return f"The client is {verb_phrase} {asset_description} in {country}."
    
    elif len(quoted_values) == 1:
        asset_description = quoted_values[0]
        verb_phrase = 'operating' if 'operat' in prompt_lower else 'developing'
        return f"The client is {verb_phrase} {asset_description}."
    
    # Fallback
    return "The client is developing renewable energy projects."


def generate_project_highlight_3para(prompt_text: str) -> str:
    """
    Generate a 3-paragraph project highlight.
    
    Structure:
    - Paragraph 1: Company operations + Project focus
    - Paragraph 2: Service scope + Partnerships + Benefits
    - Paragraph 3: Investment details
    
    Args:
        prompt_text: The prompt with substituted values
    
    Returns:
        Three paragraphs separated by double newlines
    """
    data = extract_structured_prompt_data(prompt_text)
    paragraphs = []
    
    # **PARAGRAPH 1: Company Operations + Project Focus**
    if 'paragraph_1' in data and len(data['paragraph_1']) >= 1:
        company_ops = data['paragraph_1'][0]
        project_focus = data['paragraph_1'][1] if len(data['paragraph_1']) > 1 else ""
        
        # Handle "called" pattern: extract project focus from 'called "{Project_Focus}"'
        if not project_focus and 'called' in prompt_text:
            called_match = re.search(r'called\s+["\']([^"\']+)["\']', prompt_text, re.IGNORECASE)
            if called_match:
                project_focus = called_match.group(1)
        
        sent1 = f"The client {company_ops}."
        sent2 = f"The proposed project is the {project_focus}." if project_focus else ""
        paragraphs.append(f"{sent1} {sent2}".strip())
    
    # **PARAGRAPH 2: Service Scope + Partnerships**
    if 'paragraph_2' in data and len(data['paragraph_2']) >= 1:
        service_scope = data['paragraph_2'][0]
        partnership = data['paragraph_2'][1] if len(data['paragraph_2']) > 1 else ""
        partnership_benefit = data['paragraph_2'][2] if len(data['paragraph_2']) > 2 else ""
        
        sent1 = f"The client {service_scope}."
        
        # Handle partnership verb forms
        if partnership:
            partnership_clean = partnership.strip()
            if partnership_clean.startswith('partner with') or partnership_clean.startswith('partner in'):
                sent2 = f"They {partnership_clean}."
            elif 'partnered' in partnership_clean:
                sent2 = f"They are {partnership_clean}."
            else:
                sent2 = f"They {partnership_clean}."
        else:
            sent2 = ""
        
        sent3 = f"{partnership_benefit}." if partnership_benefit else ""
        
        para2 = f"{sent1} {sent2} {sent3}".strip()
        paragraphs.append(para2)
    
    # **PARAGRAPH 3: Investment Details**
    if 'paragraph_3' in data or 'all_values' in data:
        # Look for the investment sentence pattern
        investment_pattern = re.search(
            r'The project will commence with an initial investment of\s+([^,]+),\s*with a projected expansion to\s+([^.]+)',
            prompt_text,
            re.IGNORECASE
        )
        
        if investment_pattern:
            para3 = f"The project will commence with an initial investment of {investment_pattern.group(1)}, with a projected expansion to {investment_pattern.group(2)}."
            paragraphs.append(para3)
    
    # Join paragraphs
    result = "\n\n".join(p for p in paragraphs if p)
    return result if result else "Project details are being finalized."


def extract_structured_prompt_data(prompt_text: str) -> dict:
    """
    Extract structured data from prompts that follow the pattern:
    'Paragraph 1: ... Use: "{value}"'
    
    Returns dict of extracted values and their context.
    """
    data = {}
    
    # Map them to paragraph numbers if available
    paragraphs = re.split(r'Paragraph \d+:', prompt_text)
    
    for idx, para in enumerate(paragraphs[1:], 1):  # Skip first split (before any paragraph)
        # Extract all quoted strings in this paragraph
        quoted = re.findall(r'["\']([^"\']{10,})["\']', para)  # At least 10 chars to avoid short labels
        if quoted:
            data[f'paragraph_{idx}'] = quoted
    
    # Also extract all quoted strings globally as fallback
    all_quoted = re.findall(r'["\']([^"\']{10,})["\']', prompt_text)
    data['all_values'] = all_quoted
    
    return data


def post_process_length(
    text: str,
    target_words: tuple = None,
    target_sentences: int = None,
    target_paragraphs: int = None
) -> str:
    """
    Post-process generated text to match target length requirements.
    
    Args:
        text: The generated text
        target_words: Tuple of (min, max) words
        target_sentences: Target number of sentences
        target_paragraphs: Target number of paragraphs
    
    Returns:
        Text adjusted to match length requirements
    """
    # Word count trimming
    if target_words:
        words = text.split()
        if len(words) > target_words[1]:
            text = ' '.join(words[:target_words[1]])
            # Ensure it ends with proper punctuation
            if not text.endswith('.'):
                text += '.'
    
    # Sentence count trimming
    if target_sentences:
        sentences = re.split(r'[.!?]+\s+', text)
        sentences = [s.strip() for s in sentences if s.strip()]
        if len(sentences) > target_sentences:
            text = '. '.join(sentences[:target_sentences])
            if not text.endswith('.'):
                text += '.'
    
    # Paragraph count trimming
    if target_paragraphs:
        paragraphs = text.split('\n\n')
        paragraphs = [p.strip() for p in paragraphs if p.strip()]
        if len(paragraphs) > target_paragraphs:
            text = '\n\n'.join(paragraphs[:target_paragraphs])
    
    return text.strip()


def generate_generic_content(prompt_text: str, tone: str = "short") -> str:
    """
    Fallback generic content generator for prompts that don't match fixed templates.
    
    Args:
        prompt_text: The prompt with substituted values
        tone: "short", "medium", or "long"
    
    Returns:
        Generated text
    """
    prompt_lower = prompt_text.lower()
    
    # Detect specific length requirements in the prompt
    word_count = None
    sentence_count = None
    paragraph_count = None
    
    # Look for "X words", "X-Y words"
    word_match = re.search(r'(\d+)(?:-(\d+))?\s+words?', prompt_lower)
    if word_match:
        word_count = (int(word_match.group(1)), int(word_match.group(2) or word_match.group(1)))
    
    # Look for "X sentences", "X-Y sentences"
    sent_match = re.search(r'(\d+)(?:-(\d+))?\s+sentences?', prompt_lower)
    if sent_match:
        sentence_count = (int(sent_match.group(1)), int(sent_match.group(2) or sent_match.group(1)))
    
    # Look for "X paragraphs", "X-Y paragraphs"
    para_match = re.search(r'(\d+)(?:-(\d+))?\s+paragraphs?', prompt_lower)
    if para_match:
        paragraph_count = (int(para_match.group(1)), int(para_match.group(2) or para_match.group(1)))
    
    # Extract all the actual data values from the prompt (these replaced the tokens)
    values = []
    
    # Find quoted text
    quoted = re.findall(r'"([^"]+)"', prompt_text)
    values.extend(quoted)
    
    # Find dollar amounts
    amounts = re.findall(r'\$[\d,]+(?:\s*(?:million|billion|thousand))?', prompt_text, re.IGNORECASE)
    values.extend(amounts)
    
    # Find capitalized multi-word phrases
    capitalized = re.findall(r'\b[A-Z][a-zA-Z]+(?:\s+[A-Z][a-zA-Z]+)*\b', prompt_text)
    excluded = {'Using', 'Use', 'Write', 'Follow', 'Sentence', 'Do', 'Excel', 'English', 'Initial', 'Future', 'The', 'Based'}
    capitalized = [c for c in capitalized if c not in excluded and len(c) > 2]
    values.extend(capitalized)
    
    # Find longer descriptive phrases (likely field values)
    technical_terms = re.findall(r'\b[a-z]+(?:\s+[a-z]+){1,5}\b', prompt_text)
    technical_terms = [t for t in technical_terms if len(t) > 15 and 'write' not in t and 'using' not in t and 'only' not in t]
    values.extend(technical_terms[:3])
    
    # Remove duplicates
    seen = set()
    unique_values = []
    for v in values:
        v_clean = v.strip()
        if v_clean and v_clean not in seen and len(v_clean) > 1:
            seen.add(v_clean)
            unique_values.append(v_clean)
    
    # Generate content based on requirements
    if word_count and word_count[1] <= 10:
        # Very short: 3-10 words
        if len(unique_values) >= 1:
            return unique_values[0]
        else:
            return "Professional solutions"
    
    elif paragraph_count or (sentence_count and sentence_count[1] >= 3):
        # Multi-sentence or paragraph content
        sentences = []
        target_sentences = sentence_count[1] if sentence_count else (paragraph_count[1] * 3 if paragraph_count else 3)
        
        # Build multiple sentences
        if len(unique_values) >= 2:
            sentences.append(f"{unique_values[0]} delivers innovative solutions in {unique_values[1]}.")
        elif len(unique_values) >= 1:
            sentences.append(f"{unique_values[0]} represents excellence in their field.")
        else:
            sentences.append("Professional solutions delivered with expertise.")
        
        if target_sentences >= 2:
            if len(unique_values) >= 3:
                sentences.append(f"Leveraging {unique_values[2]}, the organization drives exceptional outcomes.")
            else:
                sentences.append("Advanced capabilities ensure outstanding results.")
        
        if target_sentences >= 3:
            sentences.append("Proven track record of success across diverse initiatives.")
        
        if target_sentences >= 4:
            sentences.append("Strategic partnerships and innovation create lasting value.")
        
        return " ".join(sentences[:target_sentences])
    
    else:
        # Default: short single sentence
        if len(unique_values) >= 2:
            return f"{unique_values[0]} delivers excellence in {unique_values[1]}."
        elif len(unique_values) >= 1:
            return f"{unique_values[0]} drives innovation and results."
        else:
            return "Excellence delivered through proven expertise."


# ---------------- VALIDATION & QUALITY CHECKS ----------------

def extract_company_names(row_data: pd.Series) -> list:
    """
    Extract potential company/client names from row data.
    
    Args:
        row_data: The data row
    
    Returns:
        List of potential company names (lowercased for comparison)
    """
    potential_names = []
    
    # Common column names that contain company identifiers
    name_columns = [
        'company_name', 'company', 'client', 'client_name', 
        'issuer', 'sponsor', 'sponsor_name', 'name',
        'company legal name', 'legal name', 'project name'
    ]
    
    for col in row_data.index:
        col_lower = str(col).lower()
        if any(name_key in col_lower for name_key in name_columns):
            value = str(row_data[col])
            if value and value.lower() not in ['nan', 'none', '']:
                # Extract significant words (3+ chars) as potential name parts
                words = value.split()
                for word in words:
                    if len(word) >= 3:
                        potential_names.append(word.lower())
                # Also add the full value
                potential_names.append(value.lower())
    
    return list(set(potential_names))


def validate_anonymous_text(text: str, row_data: pd.Series) -> tuple:
    """
    Validate that text maintains anonymity (no company names leaked).
    
    Args:
        text: Generated text to validate
        row_data: Source data to check against
    
    Returns:
        Tuple of (is_valid: bool, reason: str)
    """
    text_lower = text.lower()
    
    # Extract potential company names from data
    company_names = extract_company_names(row_data)
    
    # Check if any company name appears in the text
    for name in company_names:
        if len(name) > 3 and name in text_lower:
            # Exclude common words that might be false positives
            common_words = {'the', 'client', 'company', 'project', 'group', 'inc', 'ltd', 'corp', 'llc'}
            if name not in common_words:
                return (False, f"Company name '{name}' appears in text (anonymity breach)")
    
    return (True, "")


def validate_client_oneliner(text: str, row_data: pd.Series) -> tuple:
    """
    Validate that a client one-liner meets requirements.
    
    Requirements:
    1. Must start with "The client"
    2. Must be a single sentence
    3. Must maintain anonymity (no company names)
    
    Args:
        text: Generated text to validate
        row_data: Source data to check against
    
    Returns:
        Tuple of (is_valid: bool, reason: str)
    """
    text_stripped = text.strip()
    
    # Check 1: Starts with "The client"
    if not text_stripped.startswith("The client"):
        return (False, "Does not start with 'The client'")
    
    # Check 2: Single sentence (should end with one period)
    sentences = [s.strip() for s in text_stripped.split('.') if s.strip()]
    if len(sentences) != 1:
        return (False, f"Contains {len(sentences)} sentences, expected 1")
    
    # Check 3: Anonymity
    is_anonymous, reason = validate_anonymous_text(text, row_data)
    if not is_anonymous:
        return (False, reason)
    
    return (True, "")


def make_stricter_prompt(original_prompt: str, validation_failure: str) -> str:
    """
    Create a stricter version of the prompt based on validation failure.
    
    Args:
        original_prompt: The original prompt that failed
        validation_failure: Reason why validation failed
    
    Returns:
        Stricter prompt with explicit instructions
    """
    prompt_lower = original_prompt.lower()
    
    # If it's a client one-liner that failed
    if 'the client' in prompt_lower and 'anonymous' in prompt_lower:
        # Add explicit instructions
        strict_additions = [
            "IMPORTANT: Start the sentence with exactly 'The client is'",
            "Do NOT include any company names, project names, or client identifiers",
            "Use only generic terms like 'The client', 'the project', 'the facility'",
            "Keep it to ONE sentence only"
        ]
        return original_prompt + "\n\n" + "\n".join(strict_additions)
    
    # For other types, add generic strictness
    return original_prompt + "\n\nIMPORTANT: Follow the exact format specified. Maintain complete anonymity."


def generate_beta_text(prompt_text: str, row_data: pd.Series, tone: str = "short") -> str:
    """
    Generates text by interpreting the prompt after {token} substitution.
    Uses fixed templates for common patterns with post-processing for length.
    Includes validation and regeneration logic for quality control.
    
    Args:
        prompt_text: The fully substituted prompt
        row_data: Source data (for fallback/debugging)
        tone: "short", "medium", or "long" (for fallback generation)
    
    Returns:
        Generated text matching the prompt requirements
    """
    prompt_lower = prompt_text.lower()
    
    # **TEMPLATE 1: Sector Label (2-3 words)**
    is_sector_label = (
        ('sector label' in prompt_lower or 'short label' in prompt_lower) and
        ('1-3 word' in prompt_lower or '2-3 word' in prompt_lower or 'short' in prompt_lower)
    )
    
    if is_sector_label:
        result = generate_sector_label(prompt_text, max_words=3)
        return post_process_length(result, target_words=(2, 3))
    
    # **TEMPLATE 2: Client One-liner (with validation)**
    is_client_oneliner = (
        'the client' in prompt_lower and
        ('one sentence' in prompt_lower or 'anonymous' in prompt_lower or 'briefly states' in prompt_lower)
    )
    
    if is_client_oneliner:
        # First attempt
        result = generate_client_oneliner(prompt_text)
        result = post_process_length(result, target_sentences=1)
        
        # Validate result
        is_valid, failure_reason = validate_client_oneliner(result, row_data)
        
        if not is_valid:
            # Regenerate with stricter prompt
            stricter_prompt = make_stricter_prompt(prompt_text, failure_reason)
            result = generate_client_oneliner(stricter_prompt)
            result = post_process_length(result, target_sentences=1)
            
            # Force "The client" prefix if still missing
            if not result.strip().startswith("The client"):
                result = "The client " + result.strip()
                if not result.endswith('.'):
                    result += '.'
        
        return result
    
    # **TEMPLATE 3: Project Highlight (3 paragraphs) - with anonymity check**
    is_project_highlight = (
        'paragraph 1:' in prompt_lower and 
        'paragraph 2:' in prompt_lower and 
        'paragraph 3:' in prompt_lower
    )
    
    if is_project_highlight:
        # Check if anonymity is required (look for "anonymous" or "the client" keywords)
        requires_anonymity = 'anonymous' in prompt_lower or 'the client' in prompt_lower
        
        # First attempt
        result = generate_project_highlight_3para(prompt_text)
        result = post_process_length(result, target_paragraphs=3)
        
        # Validate anonymity if required
        if requires_anonymity:
            is_anonymous, failure_reason = validate_anonymous_text(result, row_data)
            
            if not is_anonymous:
                # Regenerate with stricter prompt
                stricter_prompt = make_stricter_prompt(prompt_text, failure_reason)
                result = generate_project_highlight_3para(stricter_prompt)
                result = post_process_length(result, target_paragraphs=3)
        
        return result
    
    # **FALLBACK: Generic content generation**
    return generate_generic_content(prompt_text, tone)


# -------- Placeholder logic --------
def resolve_tag_to_value(tag: str, data_dict: dict, placeholder_mapping: dict, module_type: str = None) -> str:
    """
    Resolve a :TAG: placeholder to its value using the mapping and data dictionary.
    
    Args:
        tag: The tag name (e.g., 'COMPANY_NAME')
        data_dict: Dictionary of normalized_label -> value
        placeholder_mapping: Dictionary of TAG -> [possible label variations]
        module_type: Detected module type ("Module1", "Module2", "Module3")
    
    Returns:
        The resolved value or empty string if not found
    """
    # If we have a detected module, prioritize module-specific labels
    if module_type:
        module_mappings = get_module_field_mappings()
        if module_type in module_mappings and tag in module_mappings[module_type]:
            module_label = module_mappings[module_type][tag]
            normalized = normalize_label(module_label)
            if normalized in data_dict:
                return data_dict[normalized]
    
    # Check if this tag is in our mapping
    if tag in placeholder_mapping:
        possible_labels = placeholder_mapping[tag]
        
        # Try each possible label variation
        for label_variation in possible_labels:
            normalized = normalize_label(label_variation)
            if normalized in data_dict:
                return data_dict[normalized]
        
        # If no exact match, try fuzzy matching
        fuzzy_match = fuzzy_match_label(possible_labels[0], list(data_dict.keys()))
        if fuzzy_match and fuzzy_match in data_dict:
            return data_dict[fuzzy_match]
    
    # Try using the tag itself as a label (for backwards compatibility)
    normalized_tag = normalize_label(tag)
    if normalized_tag in data_dict:
        return data_dict[normalized_tag]
    
    # Try fuzzy match on the tag
    fuzzy_match = fuzzy_match_label(tag, list(data_dict.keys()))
    if fuzzy_match and fuzzy_match in data_dict:
        return data_dict[fuzzy_match]
    
    return ""


def process_placeholder(
    placeholder: str,
    row_data: pd.Series,
    excel_columns: list,
    beta_tone: str,
    missing_to_blank: bool,
    data_dict: dict = None,
    placeholder_mapping: dict = None,
    excel_format: str = 'column-based',
    module_type: str = None,
) -> str:
    """
    Resolve a placeholder to final text using smart local generation.

    Supports multiple formats:
    - Direct: "Title" → row_data["Title"] (column-based)
    - Tag: ":COMPANY_NAME:" → resolved via mapping (row-based)
    - AI: "AI: Write teaser about {Company}" → smart text generation
    """
    # Handle :TAG: format (new row-based system)
    tag_match = re.match(r'^:([A-Z_]+):$', placeholder.strip())
    if tag_match and data_dict is not None and placeholder_mapping is not None:
        tag = tag_match.group(1)
        return resolve_tag_to_value(tag, data_dict, placeholder_mapping, module_type)
    
    # AI placeholder - always use local generation
    if placeholder.startswith("AI:"):
        prompt_text = placeholder[3:].strip()
        
        # Handle double "AI:" at the start (common error)
        if prompt_text.startswith("AI:"):
            prompt_text = prompt_text[3:].strip()

        # Track which tokens couldn't be found
        missing_tokens = []
        
        # Replace {{ColumnName}} or {ColumnName} with actual values from Excel row
        for col in excel_columns:
            if col in row_data.index:
                value = "" if pd.isna(row_data[col]) else str(row_data[col])
                # Support both {{}} and {} formats
                prompt_text = prompt_text.replace(f"{{{{{col}}}}}", value)  # {{ColumnName}}
                prompt_text = prompt_text.replace(f"{{{col}}}", value)       # {ColumnName}
        
        # Also support {TAG} format for row-based data
        if data_dict is not None and placeholder_mapping is not None:
            # Find all {TAG} or {{TAG}} patterns
            tokens = re.findall(r'\{\{?([A-Z_]+)\}?\}?', prompt_text)
            for token in tokens:
                value = resolve_tag_to_value(token, data_dict, placeholder_mapping, module_type)
                prompt_text = prompt_text.replace(f"{{{{{token}}}}}", value)
                prompt_text = prompt_text.replace(f"{{{token}}}", value)
        
        # Fix common typos: {{TOKEN} with only one closing brace
        prompt_text = re.sub(r'\{\{([^}]+)\}(?!\})', r'{{\1}}', prompt_text)
        
        # Find remaining unreplaced tokens to report to user
        remaining_tokens = re.findall(r'\{\{([^}]+)\}\}', prompt_text)
        remaining_tokens.extend(re.findall(r'\{([^}]+)\}', prompt_text))
        
        if remaining_tokens:
            st.warning(f"⚠️ Missing Excel columns: {', '.join(set(remaining_tokens))}. Add these columns to your Excel or remove from prompt.")
        
        # Replace any remaining {{TOKEN}} or {TOKEN} with blank instead of error message
        prompt_text = re.sub(r'\{\{[^}]+\}\}', '', prompt_text)
        prompt_text = re.sub(r'\{[^}]+\}', '', prompt_text)

        st.write(f"🔍 DEBUG - After token replacement: '{prompt_text[:200]}...'")
        
        generated = generate_beta_text(prompt_text, row_data, beta_tone)
        
        st.write(f"✅ DEBUG - Generated text: '{generated[:200]}...'")
        return generated

    # Direct placeholder - check if we have row_data (column-based) or data_dict (hybrid/row-based)
    if data_dict is not None and placeholder_mapping is not None:
        # Using hybrid/row-based mode - try to resolve from data_dict
        col_name = placeholder.strip()
        normalized = normalize_label(col_name)
        
        if normalized in data_dict:
            return data_dict[normalized]
        
        # Try fuzzy match
        fuzzy_match = fuzzy_match_label(col_name, list(data_dict.keys()))
        if fuzzy_match and fuzzy_match in data_dict:
            return data_dict[fuzzy_match]
        
        return "" if missing_to_blank else f"[MISSING FIELD: {col_name}]"
    
    # Legacy column-based mode with row_data
    col_name = placeholder.strip()
    if row_data is not None and col_name in row_data.index:
        val = row_data[col_name]
        return "" if pd.isna(val) else str(val)
    else:
        return "" if missing_to_blank else f"[MISSING COLUMN: {col_name}]"


# ---------------- File uploads ----------------
col1, col2 = st.columns(2)

with col1:
    st.subheader("Step 1: Template PPTX")
    template_file = st.file_uploader(
        "Upload your PPTX template (with [Placeholders] and [AI: ...] prompts)",
        type="pptx",
        key="template",
    )

with col2:
    st.subheader("Step 2: Excel Data")
    excel_file = st.file_uploader(
        "Upload your Excel file (columns must match placeholders)",
        type=["xlsx", "xlsm"],
        key="excel",
        help="Supports .xlsx and .xlsm files. Macros in .xlsm files are ignored for security."
    )


# ---------------- Template inspector (optional helper) ----------------
if template_file:
    with st.expander("🔍 Inspect template text boxes (optional)"):
        try:
            prs_preview = Presentation(template_file)
            st.write(f"Slides: **{len(prs_preview.slides)}**")
            slide_index = st.selectbox(
                "Slide to inspect",
                options=range(len(prs_preview.slides)),
                format_func=lambda i: f"Slide {i+1}",
            )
            slide = prs_preview.slides[slide_index]
            st.write(f"**Shapes on Slide {slide_index + 1}:**")

            for idx, shape in enumerate(slide.shapes):
                # Show text from text shapes and table cells
                if hasattr(shape, "text_frame"):
                    st.write(f"- Shape {idx} name: `{shape.name}`")
                    st.write(f"  Text: `{shape.text[:80]}`")
                elif getattr(shape, "has_table", False):
                    st.write(f"- Table {idx} with {len(shape.table.rows)} rows")
        except Exception as e:
            st.error(f"Error inspecting template: {e}")


# ---------------- Main generation flow ----------------
if template_file and excel_file:
    st.divider()

    # Read Excel - try multi-sheet first
    try:
        # Detect file type and add info for .xlsm
        file_name = excel_file.name if hasattr(excel_file, 'name') else 'unknown'
        if file_name.endswith('.xlsm'):
            st.info("ℹ️ **XLSM file detected**: Reading data only, macros are disabled for security")
        
        # Check if multi-sheet (pandas reads .xlsm same as .xlsx, ignores macros)
        all_sheets = pd.read_excel(excel_file, sheet_name=None, engine='openpyxl')
        num_sheets = len(all_sheets)
        
        # For single sheet, use traditional flow
        if num_sheets == 1:
            df = list(all_sheets.values())[0]
            is_multi_sheet = False
        else:
            # Multi-sheet Excel
            is_multi_sheet = True
            st.info(f"📚 **Multi-sheet Excel detected:** {num_sheets} sheets found")
            
            # Show sheet names
            with st.expander("📄 Sheet names"):
                for sheet_name in all_sheets.keys():
                    st.write(f"- {sheet_name}")
    except Exception as e:
        st.error(f"Error reading Excel: {e}")
        st.stop()

    # Handle single vs multi-sheet
    if not is_multi_sheet:
        if df.empty:
            st.error("Excel file has no rows.")
            st.stop()

        st.subheader("📋 Excel preview")
        st.dataframe(df, use_container_width=True)

        # Detect Excel format
        excel_format = detect_excel_format(df)
        st.info(f"📋 **Detected Excel format:** {excel_format}")
    else:
        # Multi-sheet: show first sheet as preview
        first_sheet_name = list(all_sheets.keys())[0]
        first_sheet_df = all_sheets[first_sheet_name]
        
        st.subheader(f"📋 Excel preview (showing first sheet: '{first_sheet_name}')")
        st.dataframe(first_sheet_df, use_container_width=True)
        
        st.info(f"📋 **Hybrid mode**: Reading both column headers and row-based Label/Value pairs from all sheets")
    
    # Parse data based on sheet count
    if is_multi_sheet:
        # Use multi-sheet parser
        use_namespacing = st.checkbox(
            "Use sheet namespacing",
            value=False,
            help="Prefix field names with sheet name (e.g., 'financials.total_cost')"
        )
        
        data_dict, sheet_info = parse_multi_sheet_excel(excel_file, use_namespacing)
        placeholder_mapping = create_placeholder_mapping()
        
        st.success(f"✅ Parsed {len(data_dict)} total fields from {num_sheets} sheets")
        
        with st.expander("🔍 View parsed data by sheet"):
            for info in sheet_info:
                st.write(f"**{info['name']}**: {info['format']} - {info['fields']} fields")
        
        with st.expander("🔍 View all parsed fields"):
            for key, value in list(data_dict.items())[:20]:
                st.write(f"**{key}**: {value[:100] if value else '(empty)'}...")
            if len(data_dict) > 20:
                st.write(f"... and {len(data_dict) - 20} more fields")
        
        # For multi-sheet, we don't need row selection
        row_data = None
        excel_columns = []
        excel_format = 'hybrid'  # Track that we're using hybrid mode
        
    else:
        # Single sheet: use hybrid parser (reads both column and row formats)
        data_dict = parse_hybrid_excel(df)
        placeholder_mapping = create_placeholder_mapping()
        excel_format = 'hybrid'  # Track that we're using hybrid mode
        
        st.info(f"📋 **Hybrid mode**: Reading both column headers and row-based Label/Value pairs")
        st.success(f"✅ Parsed {len(data_dict)} fields from Excel")
        
        with st.expander("🔍 View parsed fields"):
            for key, value in list(data_dict.items())[:20]:
                st.write(f"**{key}**: {value[:100] if value else '(empty)'}...")
            if len(data_dict) > 20:
                st.write(f"... and {len(data_dict) - 20} more fields")
        
        # For hybrid, we don't need to select a row
        row_data = None
        excel_columns = []

    # Load PPTX template
    try:
        prs = Presentation(template_file)
    except Exception as e:
        st.error(f"Error loading PPTX template: {e}")
        st.stop()

    st.divider()
    st.subheader("⚙️ Generation settings")

    # Info about hybrid mode
    if is_multi_sheet:
        st.info(f"ℹ️ Hybrid mode: Using unified data from all {num_sheets} sheets")
    else:
        st.info("ℹ️ Hybrid mode: Using both column headers and row-based data")

    # Button to generate
    if st.button("🚀 Generate PPTX", type="primary"):
        try:
            # Detect module type from data
            module_type = detect_module_type(data_dict)
            if module_type:
                st.success(f"✅ Using {module_type} field mappings")
            
            # Count targets: text boxes + table cells
            def count_text_targets(pres: Presentation) -> int:
                total = 0
                for s in pres.slides:
                    for shp in s.shapes:
                        if hasattr(shp, "text_frame"):
                            total += 1
                        elif getattr(shp, "has_table", False):
                            tbl = shp.table
                            for r in tbl.rows:
                                total += len(r.cells)
                return total

            total_targets = count_text_targets(prs)
            processed = 0
            progress = st.progress(0.0)
            status = st.empty()

            for slide in prs.slides:
                for shape in slide.shapes:
                    # Debug: show what type of shape we're looking at
                    st.write(f"🔍 Found shape: {shape.name} (has text_frame: {hasattr(shape, 'text_frame')}, has_table: {getattr(shape, 'has_table', False)})")
                    
                    # Skip shapes without text capability
                    if not hasattr(shape, "text_frame") and not getattr(shape, "has_table", False):
                        processed += 1
                        progress.progress(processed / max(total_targets, 1))
                        continue
                    
                    # Handle tables
                    if getattr(shape, "has_table", False):
                        table = shape.table
                        for r in table.rows:
                            for cell in r.cells:
                                original_text = cell.text
                                if not original_text:
                                    processed += 1
                                    progress.progress(processed / max(total_targets, 1))
                                    continue

                                ai_matches = re.findall(r"\[AI:\s*(.+?)\]", original_text)
                                direct_matches = re.findall(r"\[(?!AI:)(.+?)\]", original_text)
                                # Also find :TAG: format placeholders
                                tag_matches = re.findall(r":([A-Z_]+):", original_text)
                                direct_matches.extend([f":{tag}:" for tag in tag_matches])

                                if not direct_matches and not ai_matches:
                                    processed += 1
                                    progress.progress(processed / max(total_targets, 1))
                                    continue

                                new_text = original_text

                                for placeholder in direct_matches:
                                    replacement = process_placeholder(
                                        placeholder,
                                        row_data,
                                        excel_columns,
                                        beta_tone,
                                        missing_to_blank,
                                        data_dict,
                                        placeholder_mapping,
                                        excel_format,
                                        module_type,
                                    )
                                    new_text = new_text.replace(f"[{placeholder}]", replacement)
                                    # Also handle :TAG: format
                                    if placeholder.startswith(':') and placeholder.endswith(':'):
                                        new_text = new_text.replace(placeholder, replacement)

                                for ai_prompt in ai_matches:
                                    status.text(f"🤖 Generating text: {ai_prompt[:60]}...")
                                    replacement = process_placeholder(
                                        f"AI: {ai_prompt}",
                                        row_data,
                                        excel_columns,
                                        beta_tone,
                                        missing_to_blank,
                                        data_dict,
                                        placeholder_mapping,
                                        excel_format,
                                        module_type,
                                    )
                                    st.write(f"DEBUG: Original=[AI: {ai_prompt}], Replacement={replacement[:100]}")
                                    new_text = new_text.replace(f"[AI: {ai_prompt}]", replacement)

                                if new_text != original_text:
                                    tf = cell.text_frame
                                    if not tf.paragraphs:
                                        cell.text = new_text
                                    else:
                                        para = tf.paragraphs[0]
                                        if para.runs:
                                            run0 = para.runs[0]
                                            font = run0.font

                                            tf.clear()
                                            new_para = tf.paragraphs[0]
                                            new_run = new_para.add_run()
                                            new_run.text = new_text

                                            if font.size:
                                                new_run.font.size = font.size
                                            if font.name:
                                                new_run.font.name = font.name
                                            if font.bold is not None:
                                                new_run.font.bold = font.bold
                                            if font.italic is not None:
                                                new_run.font.italic = font.italic
                                            if font.color and hasattr(font.color, 'rgb'):
                                                new_run.font.color.rgb = font.color.rgb

                                            new_para.alignment = para.alignment
                                            new_para.level = para.level
                                        else:
                                            cell.text = new_text

                                processed += 1
                                progress.progress(processed / max(total_targets, 1))

                        # Move to next shape
                        continue

                    # Handle regular text shapes (text boxes, rectangles, etc.)
                    # Most shapes with text will have a text_frame attribute
                    if not hasattr(shape, "text_frame"):
                        processed += 1
                        progress.progress(processed / max(total_targets, 1))
                        continue

                    original_text = shape.text
                    if not original_text:
                        processed += 1
                        progress.progress(processed / max(total_targets, 1))
                        continue

                    st.write(f"  📝 Shape text preview: '{original_text[:100]}...'")
                    
                    # More forgiving regex that handles multi-line and missing closing brackets
                    # First try standard format: [AI: ... ]
                    ai_matches = re.findall(r"\[AI:\s*(.+?)\]", original_text, re.DOTALL)
                    
                    # If no matches, try without closing bracket (for malformed placeholders)
                    if not ai_matches:
                        ai_no_close = re.findall(r"\[AI:\s*(.+)$", original_text, re.DOTALL)
                        if ai_no_close:
                            st.warning(f"⚠️ Found AI placeholder without closing bracket ]")
                            ai_matches = ai_no_close
                    
                    direct_matches = re.findall(r"\[(?!AI:)(.+?)\]", original_text)
                    # Also find :TAG: format placeholders
                    tag_matches = re.findall(r":([A-Z_]+):", original_text)
                    direct_matches.extend([f":{tag}:" for tag in tag_matches])

                    st.write(f"  🎯 Found {len(ai_matches)} AI placeholders, {len(direct_matches)} direct placeholders")

                    if not direct_matches and not ai_matches:
                        processed += 1
                        progress.progress(processed / max(total_targets, 1))
                        continue

                    new_text = original_text

                    for placeholder in direct_matches:
                        replacement = process_placeholder(
                            placeholder,
                            row_data,
                            excel_columns,
                            beta_tone,
                            missing_to_blank,
                            data_dict,
                            placeholder_mapping,
                            excel_format,
                            module_type,
                        )
                        new_text = new_text.replace(f"[{placeholder}]", replacement)
                        # Also handle :TAG: format
                        if placeholder.startswith(':') and placeholder.endswith(':'):
                            new_text = new_text.replace(placeholder, replacement)

                    for ai_prompt in ai_matches:
                        status.text(f"🤖 Generating text: {ai_prompt[:60]}...")
                        replacement = process_placeholder(
                            f"AI: {ai_prompt}",
                            row_data,
                            excel_columns,
                            beta_tone,
                            missing_to_blank,
                            data_dict,
                            placeholder_mapping,
                            excel_format,
                            module_type,
                        )
                        st.write(f"DEBUG: Original=[AI: {ai_prompt}], Replacement={replacement[:100]}")
                        new_text = new_text.replace(f"[AI: {ai_prompt}]", replacement)

                    if new_text != original_text:
                        text_frame = shape.text_frame
                        if not text_frame.paragraphs:
                            shape.text = new_text
                        else:
                            para = text_frame.paragraphs[0]
                            if para.runs:
                                run0 = para.runs[0]
                                font = run0.font

                                text_frame.clear()
                                new_para = text_frame.paragraphs[0]
                                new_run = new_para.add_run()
                                new_run.text = new_text

                                if font.size:
                                    new_run.font.size = font.size
                                if font.name:
                                    new_run.font.name = font.name
                                if font.bold is not None:
                                    new_run.font.bold = font.bold
                                if font.italic is not None:
                                    new_run.font.italic = font.italic
                                if font.color and hasattr(font.color, 'rgb'):
                                    new_run.font.color.rgb = font.color.rgb

                                new_para.alignment = para.alignment
                                new_para.level = para.level
                            else:
                                shape.text = new_text

                    processed += 1
                    progress.progress(processed / max(total_targets, 1))

            progress.progress(1.0)
            status.text("✅ Done")

            output = BytesIO()
            prs.save(output)
            output.seek(0)

            st.success("✅ PPTX generated successfully!")
            st.download_button(
                label="📥 Download filled PPTX",
                data=output.getvalue(),
                file_name="teaser_filled.pptx",
                mime=(
                    "application/"
                    "vnd.openxmlformats-officedocument."
                    "presentationml.presentation"
                ),
            )

        except Exception as e:
            st.error(f"Error generating PPTX: {e}")
            st.exception(e)