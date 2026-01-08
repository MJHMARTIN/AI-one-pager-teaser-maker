from pptx import Presentation
from pptx.util import Inches, Pt

# Create a new presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Slide 1: Title slide with basic info
slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
title_box = slide1.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "Deal Overview: :ISSUER:"

# Add more fields
details_box = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(4))
details_frame = details_box.text_frame
details_frame.text = """Industry: :INDUSTRY:
Jurisdiction: :JURISDICTION:
Issuance Type: :ISSUANCE_TYPE:
Amount: :INITIAL_NOTIONAL:
Coupon: :COUPON_RATE: (:COUPON_FREQUENCY:)
Tenor: :TENOR:"""

# Slide 2: Summary slide
slide2 = prs.slides.add_slide(prs.slide_layouts[5])
summary_title = slide2.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.8))
summary_title.text_frame.text = "Executive Summary"

summary_box = slide2.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
summary_frame = summary_box.text_frame
summary_frame.text = ":CLIENT_SUMMARY:"

# Slide 3: Project Highlights
slide3 = prs.slides.add_slide(prs.slide_layouts[5])
highlight_title = slide3.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.8))
highlight_title.text_frame.text = "Key Highlights"

highlight_box = slide3.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
highlight_frame = highlight_box.text_frame
highlight_frame.text = ":PROJECT_HIGHLIGHT:"

# Save the presentation
prs.save('test_module_template.pptx')
print("✅ Created test_module_template.pptx with standardized tags")
print("   Tags used: :ISSUER:, :INDUSTRY:, :JURISDICTION:, :ISSUANCE_TYPE:,")
print("               :INITIAL_NOTIONAL:, :COUPON_RATE:, :COUPON_FREQUENCY:,")
print("               :TENOR:, :CLIENT_SUMMARY:, :PROJECT_HIGHLIGHT:")
