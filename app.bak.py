import streamlit as st
import pandas as pd
from pptx import Presentation
from io import BytesIO
import re
import requests

# ---------------- Streamlit page config ----------------
st.set_page_config(page_title="PPTX Teaser Generator with AI", layout="wide")
st.title("📊 PPTX Teaser Generator with AI")

st.write(
    "Upload your template PPTX and Excel data. "
    "The template can contain direct placeholders like [Title] and AI prompts like "
    "[AI: Write a short teaser about {Company} in {Industry}]. "
    "Formatting (font, size, color, bold, italic) of each text box is preserved."
)

# ---------------- Sidebar: Config ----------------
st.sidebar.header("⚙️ Configuration")

# Prefer secrets in production; fallback to text input for local testing
# Support both local nested secrets and Streamlit Cloud flat secrets
deepseek_api_key = (
    st.secrets.get("DEEPSEEK_API_KEY")
    or st.secrets.get("api", {}).get("deepseek_api_key")
    or st.sidebar.text_input(
        "DeepSeek API Key",
        type="password",
        help="Store this in Streamlit secrets for production use.",
    )
)

if st.secrets.get("DEEPSEEK_API_KEY") or st.secrets.get("api", {}).get("deepseek_api_key"):
    st.sidebar.write("✅ API Key loaded from secrets")

use_beta_fallback = st.sidebar.checkbox(
    "Use Beta Local Generator when AI unavailable",
    value=True,
    help="If the AI API returns an error (e.g., 402), fill prompts locally with simple text.",
)

missing_to_blank = st.sidebar.checkbox(
    "Treat missing columns as blank",
    value=True,
    help="If a placeholder column is missing, insert blank instead of an error tag.",
)

st.sidebar.divider()

st.sidebar.info(
    """
**Template syntax**

**1. Direct placeholders (no AI)**  
- In PowerPoint:  
  `Title: [Title]`  
  `Company: [Company Name]`  
- In Excel: columns `Title`, `Company Name`

**2. AI prompts (DeepSeek)**  
- In PowerPoint:  
  `[AI: Write a 2-sentence teaser about {Company} in the {Industry} sector]`  
- In Excel: columns `Company`, `Industry`  

Text box formatting is preserved when content is replaced.
"""
)

# ---------------- DeepSeek API call ----------------
def call_deepseek(prompt: str, api_key: str) -> str:
    """
    Call DeepSeek Chat Completion API and return generated text.
    API is OpenAI-compatible. See docs: https://api-docs.deepseek.com
    """
    if not api_key:
        return "[ERROR: No DeepSeek API key provided]"

    try:
        response = requests.post(
            "https://api.deepseek.com/chat/completions",
            headers={
                "Authorization": f"Bearer {api_key}",
                "Content-Type": "application/json",
            },
            json={
                "model": "deepseek-chat",
                "messages": [{"role": "user", "content": prompt}],
                "temperature": 0.7,
                "max_tokens": 400,
            },
            timeout=30,
        )
        if response.status_code != 200:
            return f"[ERROR: DeepSeek returned {response.status_code}]"
        data = response.json()
        return data["choices"][0]["message"]["content"].strip()
    except Exception as e:
        return f"[ERROR calling DeepSeek: {e}]"


# ---------------- Beta/local text generation ----------------
def generate_beta_text(prompt_text: str, row_data: pd.Series) -> str:
    """
    Simple local generator used when AI is unavailable.
    Creates 1-2 short sentences using common fields if present.
    """

    def get_first(*names):
        for n in names:
            if n in row_data.index and not pd.isna(row_data[n]):
                return str(row_data[n])
        return None

    company = get_first("Company", "Company Name", "Client", "Account")
    industry = get_first("Industry", "Sector", "Vertical")
    title = get_first("Title", "Headline")

    parts = []
    if prompt_text:
        parts.append(prompt_text.strip())
    if company and industry:
        parts.append(f"{company} operates in the {industry} sector.")
    elif company:
        parts.append(f"{company} operates in its sector.")
    elif industry:
        parts.append(f"Operating in the {industry} space.")
    elif title:
        parts.append(f"{title} — key highlights and value propositions.")
    else:
        parts.append("Short teaser based on provided data.")

    return " ".join(parts[:2])


# ---------------- Placeholder logic ----------------
def process_placeholder(
    placeholder: str,
    row_data: pd.Series,
    api_key: str,
    excel_columns: list,
    use_beta: bool,
    missing_to_blank: bool,
) -> str:
    """
    Resolve a placeholder to final text.

    - Direct: "Title" → row_data["Title"]
    - AI: "AI: Write teaser about {Company}" → call DeepSeek with Excel values merged.
    """
    # AI placeholder
    if placeholder.startswith("AI:"):
        prompt_text = placeholder[3:].strip()

        # Replace {ColumnName} with actual values from Excel row
        for col in excel_columns:
            if col in row_data.index:
                value = "" if pd.isna(row_data[col]) else str(row_data[col])
                prompt_text = prompt_text.replace(f"{{{col}}}", value)

        ai_text = call_deepseek(prompt_text, api_key) if api_key else "[ERROR: No DeepSeek API key provided]"
        # Fallback to beta/local generation if requested and an error occurred
        if use_beta and ai_text.startswith("[ERROR"):
            return generate_beta_text(prompt_text, row_data)
        return ai_text

    # Direct placeholder
    col_name = placeholder.strip()
    if col_name in row_data.index:
        val = row_data[col_name]
        return "" if pd.isna(val) else str(val)
    else:
        return "" if missing_to_blank else f"[MISSING COLUMN: {col_name}]"


# ---------------- File uploads ----------------
col1, col2 = st.columns(2)

with col1:
    st.subheader("Step 1: Template PPTX")
    template_file = st.file_uploader(
        "Upload your PPTX template (with [Placeholders] and [AI: ...] prompts)",
        type="pptx",
        key="template",
    )

with col2:
    st.subheader("Step 2: Excel Data")
    excel_file = st.file_uploader(
        "Upload your Excel file (columns must match placeholders)",
        type="xlsx",
        key="excel",
    )


# ---------------- Template inspector (optional helper) ----------------
if template_file:
    with st.expander("🔍 Inspect template text boxes (optional)"):
        try:
            prs_preview = Presentation(template_file)
            st.write(f"Slides: **{len(prs_preview.slides)}**")
            slide_index = st.selectbox(
                "Slide to inspect",
                options=range(len(prs_preview.slides)),
                format_func=lambda i: f"Slide {i+1}",
            )
            slide = prs_preview.slides[slide_index]
            st.write(f"**Shapes on Slide {slide_index + 1}:**")

            for idx, shape in enumerate(slide.shapes):
                # Show text from text shapes and table cells
                if hasattr(shape, "text_frame"):
                    st.write(f"- Shape {idx} name: `{shape.name}`")
                    st.write(f"  Text: `{shape.text[:80]}`")
                elif getattr(shape, "has_table", False):
                    st.write(f"- Table {idx} with {len(shape.table.rows)} rows")
        except Exception as e:
            st.error(f"Error inspecting template: {e}")


# ---------------- Main generation flow ----------------
if template_file and excel_file:
    st.divider()

    # Read Excel
    try:
        df = pd.read_excel(excel_file)
    except Exception as e:
        st.error(f"Error reading Excel: {e}")
        st.stop()

    if df.empty:
        st.error("Excel file has no rows.")
        st.stop()

    st.subheader("📋 Excel preview")
    st.dataframe(df, use_container_width=True)

    # Load PPTX template
    try:
        prs = Presentation(template_file)
    except Exception as e:
        st.error(f"Error loading PPTX template: {e}")
        st.stop()

    st.divider()
    st.subheader("⚙️ Generation settings")

    # Choose row from Excel
    row_index = st.selectbox(
        "Select Excel row to use",
        options=range(len(df)),
        format_func=lambda i: f"Row {i+1}",
    )

    # Button to generate
    if st.button("🚀 Generate PPTX", type="primary"):
        excel_columns = df.columns.tolist()
        row_data = df.iloc[row_index]

        if not deepseek_api_key:
            st.warning(
                "⚠️ No DeepSeek API key provided. Direct placeholders will work, "
                "but [AI: ...] prompts will return error text."
            )

        try:
            # Count targets: text boxes + table cells
            def count_text_targets(pres: Presentation) -> int:
                total = 0
                for s in pres.slides:
                    for shp in s.shapes:
                        if hasattr(shp, "text_frame"):
                            total += 1
                        elif getattr(shp, "has_table", False):
                            tbl = shp.table
                            for r in tbl.rows:
                                total += len(r.cells)
                return total

            total_targets = count_text_targets(prs)
            processed = 0
            progress = st.progress(0.0)
            status = st.empty()

            for slide in prs.slides:
                for shape in slide.shapes:
                    # Handle tables
                    if getattr(shape, "has_table", False):
                        table = shape.table
                        for r in table.rows:
                            for cell in r.cells:
                                original_text = cell.text
                                if not original_text:
                                    processed += 1
                                    progress.progress(processed / max(total_targets, 1))
                                    continue

                                ai_matches = re.findall(r"\[AI:\s*(.+?)\]", original_text)
                                direct_matches = re.findall(r"\[(?!AI:)(.+?)\]", original_text)

                                if not direct_matches and not ai_matches:
                                    processed += 1
                                    progress.progress(processed / max(total_targets, 1))
                                    continue

                                new_text = original_text

                                for placeholder in direct_matches:
                                    replacement = process_placeholder(
                                        placeholder,
                                        row_data,
                                        deepseek_api_key,
                                        excel_columns,
                                        use_beta_fallback,
                                        missing_to_blank,
                                    )
                                    new_text = new_text.replace(f"[{placeholder}]", replacement)

                                for ai_prompt in ai_matches:
                                    status.text(f"🤖 Generating AI text: {ai_prompt[:60]}...")
                                    replacement = process_placeholder(
                                        f"AI: {ai_prompt}",
                                        row_data,
                                        deepseek_api_key,
                                        excel_columns,
                                        use_beta_fallback,
                                        missing_to_blank,
                                    )
                                    new_text = new_text.replace(f"[AI: {ai_prompt}]", replacement)

                                if new_text != original_text:
                                    tf = cell.text_frame
                                    if not tf.paragraphs:
                                        cell.text = new_text
                                    else:
                                        para = tf.paragraphs[0]
                                        if para.runs:
                                            run0 = para.runs[0]
                                            font = run0.font

                                            tf.clear()
                                            new_para = tf.paragraphs[0]
                                            new_run = new_para.add_run()
                                            new_run.text = new_text

                                            if font.size:
                                                new_run.font.size = font.size
                                            if font.name:
                                                new_run.font.name = font.name
                                            if font.bold is not None:
                                                new_run.font.bold = font.bold
                                            if font.italic is not None:
                                                new_run.font.italic = font.italic
                                            if font.color and hasattr(font.color, "rgb"):
                                                new_run.font.color.rgb = font.color.rgb

                                            new_para.alignment = para.alignment
                                            new_para.level = para.level
                                        else:
                                            cell.text = new_text

                                processed += 1
                                progress.progress(processed / max(total_targets, 1))

                        # Move to next shape
                        continue

                    # Handle regular text shapes
                    if not hasattr(shape, "text_frame"):
                        continue

                    original_text = shape.text
                    if not original_text:
                        processed += 1
                        progress.progress(processed / max(total_targets, 1))
                        continue

                    ai_matches = re.findall(r"\[AI:\s*(.+?)\]", original_text)
                    direct_matches = re.findall(r"\[(?!AI:)(.+?)\]", original_text)

                    if not direct_matches and not ai_matches:
                        processed += 1
                        progress.progress(processed / max(total_targets, 1))
                        continue

                    new_text = original_text

                    for placeholder in direct_matches:
                        replacement = process_placeholder(
                            placeholder,
                            row_data,
                            deepseek_api_key,
                            excel_columns,
                            use_beta_fallback,
                            missing_to_blank,
                        )
                        new_text = new_text.replace(f"[{placeholder}]", replacement)

                    for ai_prompt in ai_matches:
                        status.text(f"🤖 Generating AI text: {ai_prompt[:60]}...")
                        replacement = process_placeholder(
                            f"AI: {ai_prompt}",
                            row_data,
                            deepseek_api_key,
                            excel_columns,
                            use_beta_fallback,
                            missing_to_blank,
                        )
                        new_text = new_text.replace(f"[AI: {ai_prompt}]", replacement)

                    if new_text != original_text:
                        text_frame = shape.text_frame
                        if not text_frame.paragraphs:
                            shape.text = new_text
                        else:
                            para = text_frame.paragraphs[0]
                            if para.runs:
                                run0 = para.runs[0]
                                font = run0.font

                                text_frame.clear()
                                new_para = text_frame.paragraphs[0]
                                new_run = new_para.add_run()
                                new_run.text = new_text

                                if font.size:
                                    new_run.font.size = font.size
                                if font.name:
                                    new_run.font.name = font.name
                                if font.bold is not None:
                                    new_run.font.bold = font.bold
                                if font.italic is not None:
                                    new_run.font.italic = font.italic
                                if font.color and hasattr(font.color, "rgb"):
                                    new_run.font.color.rgb = font.color.rgb

                                new_para.alignment = para.alignment
                                new_para.level = para.level
                            else:
                                shape.text = new_text

                    processed += 1
                    progress.progress(processed / max(total_targets, 1))

            progress.progress(1.0)
            status.text("✅ Done")

            output = BytesIO()
            prs.save(output)
            output.seek(0)

            st.success("✅ PPTX generated successfully!")
            st.download_button(
                label="📥 Download filled PPTX",
                data=output.getvalue(),
                file_name="teaser_filled.pptx",
                mime=(
                    "application/"
                    "vnd.openxmlformats-officedocument."
                    "presentationml.presentation"
                ),
            )

        except Exception as e:
            st.error(f"Error generating PPTX: {e}")
            st.exception(e)
import streamlit as st
import pandas as pd
from pptx import Presentation
from io import BytesIO
import re
import requests

# ---------------- Streamlit page config ----------------
st.set_page_config(page_title="PPTX Teaser Generator with AI", layout="wide")
st.title("📊 PPTX Teaser Generator with AI")

st.write(
    "Upload your template PPTX and Excel data. "
    "The template can contain direct placeholders like [Title] and AI prompts like "
    "[AI: Write a short teaser about {Company} in {Industry}]. "
    "Formatting (font, size, color, bold, italic) of each text box is preserved."
)

# ---------------- Sidebar: Config ----------------
st.sidebar.header("⚙️ Configuration")

# Prefer secrets in production; fallback to text input for local testing
# Support both local nested secrets and Streamlit Cloud flat secrets
deepseek_api_key = (
    st.secrets.get("DEEPSEEK_API_KEY")
    or st.secrets.get("api", {}).get("deepseek_api_key")
    or st.sidebar.text_input(
        "DeepSeek API Key",
        type="password",
        help="Store this in Streamlit secrets for production use.",
    )
)

if st.secrets.get("DEEPSEEK_API_KEY") or st.secrets.get("api", {}).get("deepseek_api_key"):
    st.sidebar.write("✅ API Key loaded from secrets")

use_beta_fallback = st.sidebar.checkbox(
    "Use Beta Local Generator when AI unavailable",
    value=True,
    help="If the AI API returns an error (e.g., 402), fill prompts locally with simple text.",
)

missing_to_blank = st.sidebar.checkbox(
    "Treat missing columns as blank",
    value=True,
    help="If a placeholder column is missing, insert blank instead of an error tag.",
)

st.sidebar.divider()

st.sidebar.info(
    """
**Template syntax**

**1. Direct placeholders (no AI)**  
- In PowerPoint:  
  `Title: [Title]`  
  `Company: [Company Name]`  
- In Excel: columns `Title`, `Company Name`

**2. AI prompts (DeepSeek)**  
- In PowerPoint:  
  `[AI: Write a 2-sentence teaser about {Company} in the {Industry} sector]`  
- In Excel: columns `Company`, `Industry`  

Text box formatting is preserved when content is replaced.
"""
)

# ---------------- DeepSeek API call ----------------
def call_deepseek(prompt: str, api_key: str) -> str:
    """
    Call DeepSeek Chat Completion API and return generated text.
    API is OpenAI-compatible. See docs: https://api-docs.deepseek.com
    """
    if not api_key:
        return "[ERROR: No DeepSeek API key provided]"

    try:
        response = requests.post(
            "https://api.deepseek.com/chat/completions",
            headers={
                "Authorization": f"Bearer {api_key}",
                "Content-Type": "application/json",
            },
            json={
                "model": "deepseek-chat",
                "messages": [{"role": "user", "content": prompt}],
                "temperature": 0.7,
                "max_tokens": 400,
            },
            timeout=30,
        )
        if response.status_code != 200:
            return f"[ERROR: DeepSeek returned {response.status_code}]"
        data = response.json()
        return data["choices"][0]["message"]["content"].strip()
    except Exception as e:
        return f"[ERROR calling DeepSeek: {e}]"


# ---------------- Beta/local text generation ----------------
def generate_beta_text(prompt_text: str, row_data: pd.Series) -> str:
    """
    Simple local generator used when AI is unavailable.
    Creates 1-2 short sentences using common fields if present.
    """

    def get_first(*names):
        for n in names:
            if n in row_data.index and not pd.isna(row_data[n]):
                return str(row_data[n])
        return None

    company = get_first("Company", "Company Name", "Client", "Account")
    industry = get_first("Industry", "Sector", "Vertical")
    title = get_first("Title", "Headline")

    parts = []
    if prompt_text:
        parts.append(prompt_text.strip())
    if company and industry:
        parts.append(f"{company} operates in the {industry} sector.")
    elif company:
        parts.append(f"{company} operates in its sector.")
    elif industry:
        parts.append(f"Operating in the {industry} space.")
    elif title:
        parts.append(f"{title} — key highlights and value propositions.")
    else:
        parts.append("Short teaser based on provided data.")

    return " ".join(parts[:2])


# ---------------- Placeholder logic ----------------
def process_placeholder(
    placeholder: str,
    row_data: pd.Series,
    api_key: str,
    excel_columns: list,
    use_beta: bool,
    missing_to_blank: bool,
) -> str:
    """
    Resolve a placeholder to final text.

    - Direct: "Title" → row_data["Title"]
    - AI: "AI: Write teaser about {Company}" → call DeepSeek with Excel values merged.
    """
    # AI placeholder
    if placeholder.startswith("AI:"):
        prompt_text = placeholder[3:].strip()

        # Replace {ColumnName} with actual values from Excel row
        for col in excel_columns:
            if col in row_data.index:
                value = "" if pd.isna(row_data[col]) else str(row_data[col])
                prompt_text = prompt_text.replace(f"{{{col}}}", value)

        ai_text = call_deepseek(prompt_text, api_key) if api_key else "[ERROR: No DeepSeek API key provided]"
        # Fallback to beta/local generation if requested and an error occurred
        if use_beta and ai_text.startswith("[ERROR"):
            return generate_beta_text(prompt_text, row_data)
        return ai_text

    # Direct placeholder
    col_name = placeholder.strip()
    if col_name in row_data.index:
        val = row_data[col_name]
        return "" if pd.isna(val) else str(val)
    else:
        return "" if missing_to_blank else f"[MISSING COLUMN: {col_name}]"


# ---------------- File uploads ----------------
col1, col2 = st.columns(2)

with col1:
    st.subheader("Step 1: Template PPTX")
    template_file = st.file_uploader(
        "Upload your PPTX template (with [Placeholders] and [AI: ...] prompts)",
        type="pptx",
        key="template",
    )

with col2:
    st.subheader("Step 2: Excel Data")
    excel_file = st.file_uploader(
        "Upload your Excel file (columns must match placeholders)",
        type="xlsx",
        key="excel",
    )


# ---------------- Template inspector (optional helper) ----------------
if template_file:
    with st.expander("🔍 Inspect template text boxes (optional)"):
        try:
            prs_preview = Presentation(template_file)
            st.write(f"Slides: **{len(prs_preview.slides)}**")
            slide_index = st.selectbox(
                "Slide to inspect",
                options=range(len(prs_preview.slides)),
                format_func=lambda i: f"Slide {i+1}",
            )
            slide = prs_preview.slides[slide_index]
            st.write(f"**Shapes on Slide {slide_index + 1}:**")

            for idx, shape in enumerate(slide.shapes):
                # Show text from text shapes and table cells
                if hasattr(shape, "text_frame"):
                    st.write(f"- Shape {idx} name: `{shape.name}`")
                    st.write(f"  Text: `{shape.text[:80]}`")
                elif getattr(shape, "has_table", False):
                    st.write(f"- Table {idx} with {len(shape.table.rows)} rows")
        except Exception as e:
            st.error(f"Error inspecting template: {e}")


# ---------------- Main generation flow ----------------
if template_file and excel_file:
    st.divider()

    # Read Excel
    try:
        df = pd.read_excel(excel_file)
    except Exception as e:
        st.error(f"Error reading Excel: {e}")
        st.stop()

    if df.empty:
        st.error("Excel file has no rows.")
        st.stop()

    st.subheader("📋 Excel preview")
    st.dataframe(df, use_container_width=True)

    # Load PPTX template
    try:
        prs = Presentation(template_file)
    except Exception as e:
        st.error(f"Error loading PPTX template: {e}")
        st.stop()

    st.divider()
    st.subheader("⚙️ Generation settings")

    # Choose row from Excel
    row_index = st.selectbox(
        "Select Excel row to use",
        options=range(len(df)),
        format_func=lambda i: f"Row {i+1}",
    )

    # Button to generate
    if st.button("🚀 Generate PPTX", type="primary"):
        excel_columns = df.columns.tolist()
        row_data = df.iloc[row_index]

        if not deepseek_api_key:
            st.warning(
                "⚠️ No DeepSeek API key provided. Direct placeholders will work, "
                "but [AI: ...] prompts will return error text."
            )

        try:
            # Count targets: text boxes + table cells
            def count_text_targets(pres: Presentation) -> int:
                total = 0
                for s in pres.slides:
                    for shp in s.shapes:
                        if hasattr(shp, "text_frame"):
                            total += 1
                        elif getattr(shp, "has_table", False):
                            tbl = shp.table
                            for r in tbl.rows:
                                total += len(r.cells)
                return total

            total_targets = count_text_targets(prs)
            processed = 0
            progress = st.progress(0.0)
            status = st.empty()

            for slide in prs.slides:
                for shape in slide.shapes:
                    # Handle tables
                    if getattr(shape, "has_table", False):
                        table = shape.table
                        for r in table.rows:
                            for cell in r.cells:
                                original_text = cell.text
                                if not original_text:
                                    processed += 1
                                    progress.progress(processed / max(total_targets, 1))
                                    continue

                                ai_matches = re.findall(r"\[AI:\s*(.+?)\]", original_text)
                                direct_matches = re.findall(r"\[(?!AI:)(.+?)\]", original_text)

                                if not direct_matches and not ai_matches:
                                    processed += 1
                                    progress.progress(processed / max(total_targets, 1))
                                    continue

                                new_text = original_text

                                for placeholder in direct_matches:
                                    replacement = process_placeholder(
                                        placeholder,
                                        row_data,
                                        deepseek_api_key,
                                        excel_columns,
                                        use_beta_fallback,
                                        missing_to_blank,
                                    )
                                    new_text = new_text.replace(f"[{placeholder}]", replacement)

                                for ai_prompt in ai_matches:
                                    status.text(f"🤖 Generating AI text: {ai_prompt[:60]}...")
                                    replacement = process_placeholder(
                                        f"AI: {ai_prompt}",
                                        row_data,
                                        deepseek_api_key,
                                        excel_columns,
                                        use_beta_fallback,
                                        missing_to_blank,
                                    )
                                    new_text = new_text.replace(f"[AI: {ai_prompt}]", replacement)

                                if new_text != original_text:
                                    tf = cell.text_frame
                                    if not tf.paragraphs:
                                        cell.text = new_text
                                    else:
                                        para = tf.paragraphs[0]
                                        if para.runs:
                                            run0 = para.runs[0]
                                            font = run0.font

                                            tf.clear()
                                            new_para = tf.paragraphs[0]
                                            new_run = new_para.add_run()
                                            new_run.text = new_text

                                            if font.size:
                                                new_run.font.size = font.size
                                            if font.name:
                                                new_run.font.name = font.name
                                            if font.bold is not None:
                                                new_run.font.bold = font.bold
                                            if font.italic is not None:
                                                new_run.font.italic = font.italic
                                            if font.color and hasattr(font.color, "rgb"):
                                                new_run.font.color.rgb = font.color.rgb

                                            new_para.alignment = para.alignment
                                            new_para.level = para.level
                                        else:
                                            cell.text = new_text

                                processed += 1
                                progress.progress(processed / max(total_targets, 1))

                        # Move to next shape
                        continue

                    # Handle regular text shapes
                    if not hasattr(shape, "text_frame"):
                        continue

                    original_text = shape.text
                    if not original_text:
                        processed += 1
                        progress.progress(processed / max(total_targets, 1))
                        continue

                    ai_matches = re.findall(r"\[AI:\s*(.+?)\]", original_text)
                    direct_matches = re.findall(r"\[(?!AI:)(.+?)\]", original_text)

                    if not direct_matches and not ai_matches:
                        processed += 1
                        progress.progress(processed / max(total_targets, 1))
                        continue

                    new_text = original_text

                    for placeholder in direct_matches:
                        replacement = process_placeholder(
                            placeholder,
                            row_data,
                            deepseek_api_key,
                            excel_columns,
                            use_beta_fallback,
                            missing_to_blank,
                        )
                        new_text = new_text.replace(f"[{placeholder}]", replacement)

                    for ai_prompt in ai_matches:
                        status.text(f"🤖 Generating AI text: {ai_prompt[:60]}...")
                        replacement = process_placeholder(
                            f"AI: {ai_prompt}",
                            row_data,
                            deepseek_api_key,
                            excel_columns,
                            use_beta_fallback,
                            missing_to_blank,
                        )
                        new_text = new_text.replace(f"[AI: {ai_prompt}]", replacement)

                    if new_text != original_text:
                        text_frame = shape.text_frame
                        if not text_frame.paragraphs:
                            shape.text = new_text
                        else:
                            para = text_frame.paragraphs[0]
                            if para.runs:
                                run0 = para.runs[0]
                                font = run0.font

                                text_frame.clear()
                                new_para = text_frame.paragraphs[0]
                                new_run = new_para.add_run()
                                new_run.text = new_text

                                if font.size:
                                    new_run.font.size = font.size
                                if font.name:
                                    new_run.font.name = font.name
                                if font.bold is not None:
                                    new_run.font.bold = font.bold
                                if font.italic is not None:
                                    new_run.font.italic = font.italic
                                if font.color and hasattr(font.color, "rgb"):
                                    new_run.font.color.rgb = font.color.rgb

                                new_para.alignment = para.alignment
                                new_para.level = para.level
                            else:
                                shape.text = new_text

                    processed += 1
                    progress.progress(processed / max(total_targets, 1))

            progress.progress(1.0)
            status.text("✅ Done")

            output = BytesIO()
            prs.save(output)
            output.seek(0)

            st.success("✅ PPTX generated successfully!")
            st.download_button(
                label="📥 Download filled PPTX",
                data=output.getvalue(),
                file_name="teaser_filled.pptx",
                mime=(
                    "application/"
                    "vnd.openxmlformats-officedocument."
                    "presentationml.presentation"
                ),
            )

        except Exception as e:
            st.error(f"Error generating PPTX: {e}")
            st.exception(e)
import streamlit as st
import pandas as pd
from pptx import Presentation
from io import BytesIO
import re
import requests

# ---------------- Streamlit page config ----------------
st.set_page_config(page_title="PPTX Teaser Generator with AI", layout="wide")
st.title("📊 PPTX Teaser Generator with AI")

st.write(
    "Upload your template PPTX and Excel data. "
    "The template can contain direct placeholders like [Title] and AI prompts like "
    "[AI: Write a short teaser about {Company} in {Industry}]. "
    "Formatting (font, size, color, bold, italic) of each text box is preserved."
)

# ---------------- Sidebar: DeepSeek config ----------------
st.sidebar.header("⚙️ Configuration")

# Prefer secrets in production; fallback to text input for local testing
# Support both local nested secrets and Streamlit Cloud flat secrets
deepseek_api_key = (
    st.secrets.get("DEEPSEEK_API_KEY") or
    st.secrets.get("api", {}).get("deepseek_api_key") or
    st.sidebar.text_input(
        "DeepSeek API Key",
        type="password",
        help="Store this in Streamlit secrets for production use."
    )
)

if st.secrets.get("DEEPSEEK_API_KEY") or st.secrets.get("api", {}).get("deepseek_api_key"):
    st.sidebar.write("✅ API Key loaded from secrets")

# Beta/local fallback options
use_beta_fallback = st.sidebar.checkbox(
    "Use Beta Local Generator when AI unavailable",
    value=True,
    help="If the AI API returns an error (e.g., 402), fill prompts locally with simple text."
)
missing_to_blank = st.sidebar.checkbox(
    "Treat missing columns as blank",
    value=True,
    help="If a placeholder column is missing, insert blank instead of an error tag."
)

st.sidebar.divider()

st.sidebar.info(
    """
**Template syntax**

**1. Direct placeholders (no AI)**  
- In PowerPoint:  
  `Title: [Title]`  
  `Company: [Company Name]`  
- In Excel: columns `Title`, `Company Name`

**2. AI prompts (DeepSeek)**  
- In PowerPoint:  
  `[AI: Write a 2-sentence teaser about {Company} in the {Industry} sector]`  
    # Button to generate
    if st.button("🚀 Generate PPTX", type="primary"):
        excel_columns = df.columns.tolist()
        row_data = df.iloc[row_index]

        if any(col not in excel_columns for col in excel_columns):
            st.warning("Some Excel columns referenced may be missing.")

        if not deepseek_api_key:
            st.warning(
                "⚠️ No DeepSeek API key provided. Direct placeholders will work, "
                "but [AI: ...] prompts will return error text."
            )

        try:
            # Count targets: text boxes + table cells
            def count_text_targets(pres: Presentation) -> int:
                total = 0
                for s in pres.slides:
                    for shp in s.shapes:
                        if hasattr(shp, "text_frame"):
                            total += 1
                        elif getattr(shp, "has_table", False):
                            tbl = shp.table
                            for r in tbl.rows:
                                total += len(r.cells)
                return total

            total_targets = count_text_targets(prs)
            processed = 0
            progress = st.progress(0.0)
            status = st.empty()

            for slide_idx, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    # Handle tables
                    if getattr(shape, "has_table", False):
                        table = shape.table
                        for r in table.rows:
                            for cell in r.cells:
                                original_text = cell.text
                                if not original_text:
                                    processed += 1
                                    progress.progress(processed / max(total_targets, 1))
                                    continue

                                ai_matches = re.findall(r"\[AI:\s*(.+?)\]", original_text)
                                direct_matches = re.findall(r"\[(?!AI:)(.+?)\]", original_text)

                                if not direct_matches and not ai_matches:
                                    processed += 1
                                    progress.progress(processed / max(total_targets, 1))
                                    continue

                                new_text = original_text

                                for placeholder in direct_matches:
                                    replacement = process_placeholder(
                                        placeholder, row_data, deepseek_api_key, excel_columns, use_beta_fallback, missing_to_blank
                                    )
                                    new_text = new_text.replace(f"[{placeholder}]", replacement)

                                for ai_prompt in ai_matches:
                                    status.text(f"🤖 Generating AI text: {ai_prompt[:60]}...")
                                    replacement = process_placeholder(
                                        f"AI: {ai_prompt}", row_data, deepseek_api_key, excel_columns, use_beta_fallback, missing_to_blank
                                    )
                                    new_text = new_text.replace(f"[AI: {ai_prompt}]", replacement)

                                if new_text != original_text:
                                    tf = cell.text_frame
                                    if not tf.paragraphs:
                                        cell.text = new_text
                                    else:
                                        para = tf.paragraphs[0]
                                        if para.runs:
                                            run0 = para.runs[0]
                                            font = run0.font

                                            tf.clear()
                                            new_para = tf.paragraphs[0]
                                            new_run = new_para.add_run()
                                            new_run.text = new_text

                                            if font.size:
                                                new_run.font.size = font.size
                                            if font.name:
                                                new_run.font.name = font.name
                                            if font.bold is not None:
                                                new_run.font.bold = font.bold
                                            if font.italic is not None:
                                                new_run.font.italic = font.italic
                                            if font.color and hasattr(font.color, 'rgb'):
                                                new_run.font.color.rgb = font.color.rgb

                                            new_para.alignment = para.alignment
                                            new_para.level = para.level
                                        else:
                                            cell.text = new_text

                                processed += 1
                                progress.progress(processed / max(total_targets, 1))

                        # Move to next shape
                        continue

                    # Handle regular text shapes
                    if not hasattr(shape, "text_frame"):
                        continue

                    original_text = shape.text
                    if not original_text:
                        processed += 1
                        progress.progress(processed / max(total_targets, 1))
                        continue

                    ai_matches = re.findall(r"\[AI:\s*(.+?)\]", original_text)
                    direct_matches = re.findall(r"\[(?!AI:)(.+?)\]", original_text)

                    if not direct_matches and not ai_matches:
                        processed += 1
                        progress.progress(processed / max(total_targets, 1))
                        continue

                    new_text = original_text

                    for placeholder in direct_matches:
                        replacement = process_placeholder(
                            placeholder, row_data, deepseek_api_key, excel_columns, use_beta_fallback, missing_to_blank
                        )
                        new_text = new_text.replace(f"[{placeholder}]", replacement)

                    for ai_prompt in ai_matches:
                        status.text(f"🤖 Generating AI text: {ai_prompt[:60]}...")
                        replacement = process_placeholder(
                            f"AI: {ai_prompt}", row_data, deepseek_api_key, excel_columns, use_beta_fallback, missing_to_blank
                        )
                        new_text = new_text.replace(f"[AI: {ai_prompt}]", replacement)

                    if new_text != original_text:
                        text_frame = shape.text_frame
                        if not text_frame.paragraphs:
                            shape.text = new_text
                        else:
                            para = text_frame.paragraphs[0]
                            if para.runs:
                                run0 = para.runs[0]
                                font = run0.font

                                text_frame.clear()
                                new_para = text_frame.paragraphs[0]
                                new_run = new_para.add_run()
                                new_run.text = new_text

                                if font.size:
                                    new_run.font.size = font.size
                                if font.name:
                                    new_run.font.name = font.name
                                if font.bold is not None:
                                    new_run.font.bold = font.bold
                                if font.italic is not None:
                                    new_run.font.italic = font.italic
                                if font.color and hasattr(font.color, 'rgb'):
                                    new_run.font.color.rgb = font.color.rgb

                                new_para.alignment = para.alignment
                                new_para.level = para.level
                            else:
                                shape.text = new_text

                    processed += 1
                    progress.progress(processed / max(total_targets, 1))

            progress.progress(1.0)
            status.text("✅ Done")

            output = BytesIO()
            prs.save(output)
            output.seek(0)

            st.success("✅ PPTX generated successfully!")
            st.download_button(
                label="📥 Download filled PPTX",
                data=output.getvalue(),
                file_name="teaser_filled.pptx",
                mime=(
                    "application/"
                    "vnd.openxmlformats-officedocument."
                    "presentationml.presentation"
                ),
            )

        except Exception as e:
            st.error(f"Error generating PPTX: {e}")
            st.exception(e)
                                        for r in tbl.rows:
                                            total += len(r.cells)
                            return total

                        total_targets = count_text_targets(prs)
                        replacement = process_placeholder(
                            f"AI: {ai_prompt}", row_data, deepseek_api_key, excel_columns, use_beta_fallback, missing_to_blank
                        )
                        new_text = new_text.replace(f"[AI: {ai_prompt}]", replacement)

                    # If changed, update while preserving formatting
                                # Handle table cells
                                if getattr(shape, "has_table", False):
                                    table = shape.table
                                    for r in table.rows:
                                        for cell in r.cells:
                                            original_text = cell.text
                                            if not original_text:
                                                processed += 1
                                                progress.progress(processed / max(total_targets, 1))
                                                continue

                                            # Find placeholders
                                            ai_matches = re.findall(r"\[AI:\s*(.+?)\]", original_text)
                                            direct_matches = re.findall(r"\[(?!AI:)(.+?)\]", original_text)

                                            if not direct_matches and not ai_matches:
                                                processed += 1
                                                progress.progress(processed / max(total_targets, 1))
                                                continue

                                            new_text = original_text

                                            # Replace direct placeholders
                                            for placeholder in direct_matches:
                                                replacement = process_placeholder(
                                                    placeholder, row_data, deepseek_api_key, excel_columns, use_beta_fallback, missing_to_blank
                                                )
                                                new_text = new_text.replace(f"[{placeholder}]", replacement)

                                            # Replace AI prompts
                                            for ai_prompt in ai_matches:
                                                status.text(f"🤖 Generating AI text: {ai_prompt[:60]}...")
                                                replacement = process_placeholder(
                                                    f"AI: {ai_prompt}", row_data, deepseek_api_key, excel_columns, use_beta_fallback, missing_to_blank
                                                )
                                                new_text = new_text.replace(f"[AI: {ai_prompt}]", replacement)

                                            # If changed, update while preserving formatting
                                            if new_text != original_text:
                                                tf = cell.text_frame
                                                if not tf.paragraphs:
                                                    cell.text = new_text
                                                else:
                                                    para = tf.paragraphs[0]
                                                    if para.runs:
                                                        run0 = para.runs[0]
                                                        font = run0.font

                                                        tf.clear()
                                                        new_para = tf.paragraphs[0]
                                                        new_run = new_para.add_run()
                                                        new_run.text = new_text

                                                        if font.size:
                                                            new_run.font.size = font.size
                                                        if font.name:
                                                            new_run.font.name = font.name
                                                        if font.bold is not None:
                                                            new_run.font.bold = font.bold
                                                        if font.italic is not None:
                                                            new_run.font.italic = font.italic
                                                        if font.color and hasattr(font.color, 'rgb'):
                                                            new_run.font.color.rgb = font.color.rgb

                                                        new_para.alignment = para.alignment
                                                        new_para.level = para.level
                                                    else:
                                                        cell.text = new_text

                                            processed += 1
                                            progress.progress(processed / max(total_targets, 1))

                                    # Done with table; continue to next shape
                                    continue

                                # Handle regular text shapes
                                if not hasattr(shape, "text_frame"):
                                    continue

                                original_text = shape.text
                                if not original_text:
                                    processed += 1
                                    progress.progress(processed / max(total_targets, 1))
                                    continue

                                # Find placeholders
                                ai_matches = re.findall(r"\[AI:\s*(.+?)\]", original_text)
                                direct_matches = re.findall(r"\[(?!AI:)(.+?)\]", original_text)

                                if not direct_matches and not ai_matches:
                                    processed += 1
                                    progress.progress(processed / max(total_targets, 1))
                                    continue

                                new_text = original_text

                                # Replace direct placeholders
                                for placeholder in direct_matches:
                                    replacement = process_placeholder(
                                        placeholder, row_data, deepseek_api_key, excel_columns, use_beta_fallback, missing_to_blank
                                    )
                                    new_text = new_text.replace(f"[{placeholder}]", replacement)

                                # Replace AI prompts
                                for ai_prompt in ai_matches:
                                    status.text(f"🤖 Generating AI text: {ai_prompt[:60]}...")
                                    replacement = process_placeholder(
                                        f"AI: {ai_prompt}", row_data, deepseek_api_key, excel_columns, use_beta_fallback, missing_to_blank
                                    )
                                    new_text = new_text.replace(f"[AI: {ai_prompt}]", replacement)

                                # If changed, update while preserving formatting
                                if new_text != original_text:
                                    text_frame = shape.text_frame
                                    if not text_frame.paragraphs:
                                        shape.text = new_text
                                    else:
                                        para = text_frame.paragraphs[0]
                                        if para.runs:
                                            run0 = para.runs[0]
                                            font = run0.font

                                            # Clear frame
                                            text_frame.clear()
                                            new_para = text_frame.paragraphs[0]
                                            new_run = new_para.add_run()
                                            new_run.text = new_text

                                            # Copy basic formatting
                                            if font.size:
                                                new_run.font.size = font.size
                                            if font.name:
                                                new_run.font.name = font.name
                                            if font.bold is not None:
                                                new_run.font.bold = font.bold
                                            if font.italic is not None:
                                                new_run.font.italic = font.italic
                                            if font.color and hasattr(font.color, 'rgb'):
                                                new_run.font.color.rgb = font.color.rgb

                                            # Paragraph alignment/level
                                            new_para.alignment = para.alignment
                                            new_para.level = para.level
                                        else:
                                            # No runs, fallback
                                            shape.text = new_text

                                processed += 1
                                progress.progress(processed / max(total_targets, 1))