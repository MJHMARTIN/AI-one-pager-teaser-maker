#!/usr/bin/env python3
"""Create test examples for sector labels and client one-liners."""

from pptx import Presentation
from pptx.util import Inches, Pt
import pandas as pd

# Create test Excel data for different scenarios
test_data = {
    'Label': [
        'Sector',
        'Country',
        'Project_Type',
        'Asset_Description',
    ],
    'Value': [
        'New Energy',
        'Japan',
        'Grid-connected battery storage station',
        'grid-connected battery storage station project',
    ]
}

# Create test files for each energy type
scenarios = [
    {
        'name': 'battery_japan',
        'sector': 'New Energy',
        'country': 'Japan',
        'project_type': 'Grid-connected battery storage station',
        'asset_desc': 'grid-connected battery storage station project',
        'expected_label': 'Energy Storage',
        'expected_oneliner': 'The client is developing a grid-connected battery storage station project in Japan.'
    },
    {
        'name': 'solar_vietnam',
        'sector': 'New Energy',
        'country': 'Vietnam',
        'project_type': 'Utility-scale solar photovoltaic farm',
        'asset_desc': '500 MW utility-scale solar photovoltaic farm',
        'expected_label': 'Solar Energy',
        'expected_oneliner': 'The client is developing a 500 MW utility-scale solar photovoltaic farm in Vietnam.'
    },
    {
        'name': 'wind_india',
        'sector': 'New Energy',
        'country': 'India',
        'project_type': 'Onshore wind farm',
        'asset_desc': '250 MW onshore wind farm with 100 MWh battery storage integration',
        'expected_label': 'Wind Energy',
        'expected_oneliner': 'The client is developing a 250 MW onshore wind farm with 100 MWh battery storage integration in India.'
    }
]

# Create Excel files
for scenario in scenarios:
    data = {
        'Label': ['Sector', 'Country', 'Project_Type', 'Asset_Description'],
        'Value': [scenario['sector'], scenario['country'], scenario['project_type'], scenario['asset_desc']]
    }
    df = pd.DataFrame(data)
    df.to_excel(f"test_{scenario['name']}_labels.xlsx", index=False)
    print(f"✅ Created test_{scenario['name']}_labels.xlsx")

# Create PPT template with the new prompt patterns
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Slide 1: Sector Label Test
slide1 = prs.slides.add_slide(prs.slide_layouts[6])

title1 = slide1.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
title1.text_frame.text = "Test: 2-3 Word Sector Label"
title1.text_frame.paragraphs[0].font.size = Pt(32)
title1.text_frame.paragraphs[0].font.bold = True

prompt1 = slide1.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(2))
prompt1.text_frame.text = '[AI: Given Sector = "{Sector}" and Project_Type = "{Project_Type}", generate a 1-3 word sector label (no company name), e.g. "New Energy", "Solar Energy", "Wind Energy".]'
prompt1.text_frame.paragraphs[0].font.size = Pt(16)
prompt1.text_frame.word_wrap = True

expected1 = slide1.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1.5))
expected1.text_frame.text = 'Expected output:\n• Battery/Storage → "Energy Storage"\n• Solar → "Solar Energy"\n• Wind → "Wind Energy"'
expected1.text_frame.paragraphs[0].font.size = Pt(14)

# Slide 2: Client One-Liner Test
slide2 = prs.slides.add_slide(prs.slide_layouts[6])

title2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
title2.text_frame.text = "Test: Anonymous Client One-Liner"
title2.text_frame.paragraphs[0].font.size = Pt(32)
title2.text_frame.paragraphs[0].font.bold = True

prompt2 = slide2.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(2.5))
prompt2.text_frame.text = '[AI: Using Asset_Description = "{Asset_Description}" and Country = "{Country}", write one anonymous sentence starting with "The client" that briefly states what is being developed/operated and where. Do not mention the company name.]'
prompt2.text_frame.paragraphs[0].font.size = Pt(16)
prompt2.text_frame.word_wrap = True

expected2 = slide2.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(2))
expected2.text_frame.text = '''Expected format:
"The client is developing {asset description} in {country}."

Example:
"The client is developing a 500 MW utility-scale solar photovoltaic farm in Vietnam."'''
expected2.text_frame.paragraphs[0].font.size = Pt(14)

# Save template
prs.save('test_labels_oneliner_template.pptx')
print("✅ Created test_labels_oneliner_template.pptx")

print("\n📋 Test scenarios created:")
for scenario in scenarios:
    print(f"\n{scenario['name'].upper()}:")
    print(f"  Expected label: {scenario['expected_label']}")
    print(f"  Expected one-liner: {scenario['expected_oneliner']}")
